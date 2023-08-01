from PIL import Image
import os
import requests
import json
from pptx import Presentation
from pptx.util import Inches, Pt
import io
import pptx.util



# Declare the API base URL as a global constant
API_BASE_URL = "https://abcd2.projectabcd.com/api/getinfo.php?id="

def fetch_data_from_api(ids):
    data_list = []
    for id in ids:
        url = API_BASE_URL + str(id)
        headers = {
            "Accept": "application/json",
            "User-Agent": "Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/58.0.3029.110 Safari/537.36"
        }

        response = requests.get(url, headers=headers)

        if response.status_code == 200:
            data_list.append(response.json()['data'])
        else:
            print(f"Failed to get data from the API for ID {id}. Status code:", response.status_code)

    return data_list

def parse_slide_numbers(filename):
    with open(filename, 'r') as file:
        slide_numbers = file.read()

    # Split slide numbers using comma or dash as separators
    ids = []
    for item in slide_numbers.split(','):
        if '-' in item:
            start, end = map(int, item.split('-'))
            ids.extend(range(start, end + 1))
        else:
            ids.append(int(item))

    return ids


def read_preferences(filename):
    preferences = {}
    with open(filename, 'r') as file:
        for line in file:
            key, value = line.strip().split(' = ')
            key = key.strip().lower()  # Convert the key to lowercase
            if value.isdigit():
                value = int(value)
            elif value.startswith('"') and value.endswith('"'):
                value = value[1:-1]
            preferences[key] = value
            
    return preferences


def find_content_placeholder_index(slide_layout):
    for idx, shape in enumerate(slide_layout.placeholders):
        if shape.placeholder_format.idx == 1:
            # Placeholder at index 1 corresponds to the content placeholder in slide layout 5.
            return idx
    return None

def create_pptx_and_save(fetched_data_list, preferences):
    # Read preferences
    

    # Initialize the presentation
    prs = Presentation()

    # Loop through each data and add slides accordingly
    for data in fetched_data_list:
        # Add a content slide (layout 5) to display the JSON data
        content_slide_layout = prs.slide_layouts[3]
        content_slide = prs.slides.add_slide(content_slide_layout)
       
                    
        prs.slide_width = pptx.util.Inches(11)
        prs.slide_height = pptx.util.Inches(8) 
        
                
        # Set the name as the title at the top of the slide
        title_shape = content_slide.shapes.title
        title_shape.text = data['name']

        # Find the index of the content placeholder in the layout
        content_placeholder_idx = find_content_placeholder_index(content_slide_layout)
        if content_placeholder_idx is not None:
            content_box = content_slide.placeholders[content_placeholder_idx]

             # Add the description and did_you_know to the content box text
            description_text = f"Description: {data['description']}\n"
            did_you_know_text = f"Fun Fact {data['did_you_know']}"
            content_box.text = description_text + did_you_know_text

            # Adjust font size of the content text
            text_frame = content_box.text_frame
            for paragraph in text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.size = Pt(18)  # Set the font size to 12 points (adjust as needed)
                    
            # Adjust the width of the content box (text frame)
            content_box_width = Inches(6)  # Set the width of the content box to 6 inches (adjust as needed)
            text_frame.width = content_box_width
            

       # Load the image from the file and add it to the slide
        if 'image_url' in data and data['image_url']:
            image_path = str(data['image_url'])
            image_path1 = os.path.basename(image_path)

            left_inch = Inches(4.95)  # Left position of the image
            top_inch = Inches(1.25)  # Top position of the image
            width_inch = Inches(5.5)  # Width of the image
            height_inch = Inches(6.25)  # Height of the image

            # Add the image to the slide using the Image class
            content_slide.shapes.add_picture(str(image_path1), left_inch, top_inch, width_inch, height_inch)

            
    # Save the PowerPoint presentation to a file
    output_filename = "api.pptx"
    prs.save(output_filename)

    print(f"PowerPoint presentation '{output_filename}' created successfully.")
    

if __name__ == "__main__":
    slide_numbers_file = "slide_numbers.txt"
    preferences_file = "preferences.txt"

    ids = parse_slide_numbers(slide_numbers_file)
    preferences = read_preferences(preferences_file)
    
    
    if ids:
        fetched_data_list = fetch_data_from_api(ids)
        if fetched_data_list:
            create_pptx_and_save(fetched_data_list, preferences)
        else:
            print("Failed to fetch data from the API.")
    else:
        print("No valid slide numbers found in the file.")
