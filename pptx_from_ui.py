import tkinter as tk
from tkinter import*
import os
import shutil
from os.path import basename
from collections.abc import Container
import pptx.util
from pptx import Presentation 
from pptx.util import Pt
from bs4 import BeautifulSoup
import requests
import pandas as pd
from sqlalchemy import create_engine
import mysql.connector
import sys



window = tk.Tk()
window.title("ABCD Book Generation UI")

label_text = tk.StringVar()
label_text.set("Enter page numbers")
text_box =tk.Label(window, textvariable=label_text).grid(row=1,column=1)


text_value = tk.Text(window,height=5, width= 20)
text_value.grid(row=1,column=2)
filevalue = "26, 27, 28, 29, 30, 39, 50, 52, 53, 110, 111, 112, 116, 196, 206, 262, 265, 275, 276, 314, 317, 318, 319, 320, 321, 322, 324, 325, 326, 327, 328, 329, 401, 405, 406, 407, 409, 410, 411, 412, 413, 415, 418, 419, 422, 423, 424, 425, 426, 427, 428, 431, 432, 433, 434, 437, 438, 439, 440, 441, 442, 443, 444, 445, 462, 463, 468, 469, 470, 471, 472, 475, 476, 477, 478, 483, 484, 491, 492, 493, 502, 506, 520, 542, 544, 549, 568, 574, 578, 581, 582, 601, 605, 611, 626, 627, 631, 649, 654, 655, 658, 659, 660, 662, 664, 665, 666, 667, 670, 671, 672, 674, 678, 681, 682, 683, 684, 686, 688, 689, 690, 691, 693, 695, 696, 697, 698, 699, 700"
text_value.insert("1.0", filevalue)
#label for preferences
label_text_1 = tk.StringVar()
label_text_1.set("preferences")
text_box_1 =tk.Label(window, textvariable=label_text_1).grid(row=5,column=1)

#input for preferences
nums_1 = StringVar(value="preferences.txt")
text_value_1 = tk.Entry(window, textvariable=nums_1).grid(row=5,column=2)

#label for output file
label_text_2 = tk.StringVar()
label_text_2.set("Output file")
text_box_2 =tk.Label(window, textvariable=label_text_2).grid(row=6,column=1)

#inputfor output file
nums_2 = StringVar(value="abcd_book")
text_value_2 = Entry(window,textvariable=nums_2).grid(row=6,column=2)
btn_1 = tk.Button(window, text= ".pptx").grid(row=6,column=2, columnspan=3)

#selected option variables
selected_option = tk.StringVar()
selected_option_1 = tk.StringVar()
selected_option_btn = tk.StringVar()

#label for layout
btn_text = tk.StringVar()
btn_text.set("Layout")
btn_box =tk.Label(window, textvariable=btn_text).grid(row=2,column=1)

# Create the radio buttons
radio_button5 = tk.Radiobutton(window, text="Landscape", variable = selected_option, value=5)
radio_button4 = tk.Radiobutton(window, text="pic seperate", variable = selected_option, value=4)
radio_button1 = tk.Radiobutton(window, text="Pic on left", variable=selected_option, value=3)
radio_button2 = tk.Radiobutton(window, text="Pic on right", variable=selected_option, value=2)
radio_button3 = tk.Radiobutton(window, text="Pic on top", variable=selected_option, value=1)

# Set the default selection to "Pic on left"
selected_option.set(3)
selected_option_btn.set(3)
radio_button1.grid(row=2, column=2)
radio_button2.grid(row=2, column=3)
radio_button3.grid(row=2, column=4)
radio_button4.grid(row = 2, column=5)
radio_button5.grid(row = 2, column=6)


#label for sort order
sort_btn = tk.StringVar()
sort_btn.set("Sort Order")
sort_box =tk.Label(window, textvariable=sort_btn).grid(row=4,column=1)

#sort order buttons
sort_button1 = tk.Radiobutton(window, text="By Name", variable=selected_option_btn, value=3)
sort_button2 = tk.Radiobutton(window, text="By ID", variable=selected_option_btn, value=2)
sort_button3 = tk.Radiobutton(window, text="By Input Order", variable=selected_option_btn, value=1)

# Layout the radio buttons
sort_button1.grid(row=4, column=2)
sort_button2.grid(row=4, column=3)
sort_button3.grid(row=4, column=4)

#label for Method
btn_text_1 = tk.StringVar()
btn_text_1.set("Method")
btn_box_1 =tk.Label(window, textvariable=btn_text_1).grid(row=3,column=1)

# Create the second radio button
radio_button2 = tk.Radiobutton(window, text="Web", variable=selected_option_1, value="Web")
radio_button3 = tk.Radiobutton(window, text="Excel", variable=selected_option_1, value="Excel")
radio_button4 = tk.Radiobutton(window, text="API", variable=selected_option_1, value="API")
radio_button5 = tk.Radiobutton(window, text="Database", variable=selected_option_1, value="Database")

# Set the default selection to "Web"
selected_option_1.set("Web")

# Layout the radio buttons
radio_button2.grid(row=3, column=2)
radio_button3.grid(row=3, column=3)
radio_button4.grid(row=3, column=4)
radio_button5.grid(row=3, column=5)

#label for download
download_option = tk.StringVar()
down_button = tk.Radiobutton(window, text="Download", variable=download_option, value="yes")
down_button.grid(row=1, column=3)
download_option.set(None)

def generate_output():
   
    page_numbers = text_value.get("1.0", "end").strip()
    preferences = nums_1.get()
    output_file = nums_2.get()
    layout_option = selected_option.get()
    method_option = selected_option_1.get()
    sort_option = selected_option_btn.get()
    downloading = download_option.get()

    global methods,pages,layout,preference,output,sort_choice,downlaod_choice
    methods= method_option
    pages = page_numbers
    layout = layout_option
    preference = preferences
    output = output_file
    sort_choice = sort_option
    downlaod_choice = downloading

    window.destroy()

    
btn = tk.Button(window, text="Generate", command=generate_output)
btn.grid(row=7, column=2, columnspan=2)


window.mainloop()



all_pages = []  
with open("slide_numbers.txt" , "w") as file:
   chosen_nums = pages
   for word in chosen_nums:
       file.write(word + "\n")
file.close

with open("slide_numbers.txt" , "r") as f:
    tokens = chosen_nums.split(",")
    for elem in tokens:
        if '-' in elem:
            token_parts = elem.split("-")
            first_token = int(token_parts[0])
            last_token = int(token_parts[1])
            for x in range(first_token, last_token + 1):
                all_pages.append(x)
        else:
            all_pages.append(elem)
        

if(downlaod_choice == "yes"):
    if(methods == "Web"):
        start = 1
        end = 700
        presentationLength = end - start + 1
        print("Downloading all 700 images...please wait.")
        for i in range (0, presentationLength):
                URL = "https://projectabcd.com/display_the_dress.php?id=" + str(start + i)
                page = requests.get(URL, headers={"User-Agent": "html"})
                soup = BeautifulSoup(page.content, "html.parser")
                logo = soup.find("img")
                printLogo = logo.attrs["src"]
                logoURL = "http://projectabcd.com/" + printLogo
                r = requests.get(logoURL, headers={"User-Agent": "html"}, stream=True)
                if r.status_code == 200:
                    with open(basename(printLogo), "wb") as f:
                        r.raw.decode_content = True
                        shutil.copyfileobj(r.raw, f)
                pageInfo = soup.find("div", class_="containerTitle")
                pageInfoImg = soup.find("div", class_="container")
                name = pageInfo.find("h2", class_="headTwo")
                printName = name.text
                image = pageInfoImg.find("div", class_="containerImage")
                img = image.find("image",class_= "image")
                printImage = img.get("src")
                pictureURL = "http://projectabcd.com/" + printImage
                r = requests.get(pictureURL, headers={"User-Agent": "html"}, stream=True)
                if r.status_code == 200:
                    with open(basename(printImage), "wb") as f:
                        r.raw.decode_content = True
                        shutil.copyfileobj(r.raw, f)
                element = pageInfo.p
                pagetext = pageInfoImg.find("div", class_="containerText")
                description = pagetext.find("p", class_="words")
                printDescription = description.text
                fact = description.find_next_sibling("p")
                printFact = fact.text
                printWords = printDescription + printFact
        print("Download complete, bye")
        sys.exit()
    elif(methods != "Web"):
        print("ERROR! Must be using web method to download!")
        sys.exit()


if(methods == "Web"):
    
    def buildPresentation():  
        with open(preference, "r") as f:
            slideOption = int(layout)
            textFont = f.readline().split("= ")
            textFont = textFont[1]
            titleFont = f.readline().split("= ")
            titleFont = titleFont[1]
            textSize = f.readline().split("= ")
            textSize = int(textSize[1])
            titleSize = f.readline().split("= ")
            titleSize = int(titleSize[1])
        prs = Presentation()
        
        pictureSlide = 0
        
       
        if(slideOption == 2):
            pictureSlide = 1
        #web scrapes the URL to get all needed information from the page
        presentationLength = len(all_pages)
        list_values_ = []

        if(sort_choice == "3"):
            for i in range (len(all_pages)):
                new_val = int(all_pages[i])
                list_values_.append(new_val)
            list_values_.sort()
            

        if(sort_choice == "2"):
            for i in range (len(all_pages)):
                new_val = int(all_pages[i])
                list_values_.append(new_val)
                
            list_values_.sort()
             
            
        if(sort_choice == "1"): 
            for i in range (len(all_pages)):
                new_val = int(all_pages[i])                  
                list_values_.append(new_val)
                
        for i in range (presentationLength):   
            if(sort_choice == "3"):
                value = int(list_values_[i])
                
            elif(sort_choice == "2"):
                value = int(list_values_[i])
                
            elif(sort_choice == "1"):
                value = int(all_pages[i])
                
            else:
                print("Choose a sort type")
                sys.exit()
            if(sort_choice == "2"):
                URL = "https://projectabcd.com/display_the_dress.php?id=" + str(list_values_[i])
            elif(sort_choice == "3"):
                URL = "https://projectabcd.com/display_the_dress.php?id=" + str(list_values_[i])
            elif(sort_choice == "1"):
                URL = "https://projectabcd.com/display_the_dress.php?id=" + str(list_values_[i])
            page = requests.get(URL, headers={"User-Agent": "html"})
            soup = BeautifulSoup(page.content, "html.parser")
            logo = soup.find("img")
            printLogo = logo.attrs["src"]
            logoURL = "http://projectabcd.com/" + printLogo
            r = requests.get(logoURL, headers={"User-Agent": "html"}, stream=True)
            if r.status_code == 200:
                   with open(basename(printLogo), "wb") as f:
                        r.raw.decode_content = True
                        shutil.copyfileobj(r.raw, f)
            pageInfo = soup.find("div", class_="containerTitle")
            pageInfoImg = soup.find("div", class_="container")
            name = pageInfo.find("h2", class_="headTwo")
            printName = name.text
            image = pageInfoImg.find("div", class_="containerImage")
            img = image.find("image",class_= "image")
            printImage = img.get("src")
            pictureURL = "http://projectabcd.com/" + printImage
            r = requests.get(pictureURL, headers={"User-Agent": "html"}, stream=True)
            if r.status_code == 200:
                with open(basename(printImage), "wb") as f:
                    r.raw.decode_content = True
                    shutil.copyfileobj(r.raw, f)
            
            pagetext = pageInfoImg.find("div", class_="containerText")
            description = pagetext.find("p", class_="words")
            printDescription = description.text
            fact = description.find_next_sibling("p")
            printFact = fact.text
             

            #creates the slide presentation if slide option 1 is choosen
            if(slideOption == 1 ):
                #creates the slides and sets layout preferences
                slide_layout = prs.slide_layouts[6]
                slide = prs.slides.add_slide(slide_layout)
                prs.slide_width = pptx.util.Inches(8)
                prs.slide_height = pptx.util.Inches(11)
                #places the logo on the slide
                logoHolder = slide.shapes.add_picture(basename(printLogo), pptx.util.Inches(7), pptx.util.Inches(10),width=pptx.util.Inches(1), height=pptx.util.Inches(1))
                #places the title on the slide
                titleBox = slide.shapes.add_textbox(pptx.util.Inches(2.25), pptx.util.Inches(.5), width = pptx.util.Inches(3), height = pptx.util.Inches(1))
                titleBoxtf = titleBox.text_frame
                title = titleBoxtf.add_paragraph()
                title.text = printName
                title.font.name = titleFont
                title.font.size = Pt(titleSize)
                
                slideId = slide.shapes.add_textbox(pptx.util.Inches(6.5), pptx.util.Inches(10.25), width = pptx.util.Inches(1), height = pptx.util.Inches(1))                       
                slideIdf = slideId.text_frame
                slide_id =slideIdf.add_paragraph()
                slide_id.text =str(value)
                #places the picture on the slide
                pictureHolder = prs.slides[i].shapes
                pictureHolder.add_picture(basename(printImage), pptx.util.Inches(2.5),pptx.util.Inches(2),width = pptx.util.Inches(3), height = pptx.util.Inches(4))
                #creates a textbox for the description and fun fact
                contentBox = slide.shapes.add_textbox(pptx.util.Inches(1), pptx.util.Inches(6), width = pptx.util.Inches(6), height = pptx.util.Inches(5))
                contentBoxtf = contentBox.text_frame
                contentBoxtf.word_wrap = True
                descriptionTitle = contentBoxtf.add_paragraph()
                descriptionTitle.font.name = textFont
                descriptionTitle.font.bold = True
                descriptionTitle.font.size = Pt(textSize)
                descriptionTitle.text = "Description: "
                descriptionParagraph = contentBoxtf.add_paragraph()
                descriptionParagraph.font.name = textFont
                descriptionParagraph.font.size = Pt(textSize)
                descriptionParagraph.text = printDescription
                FunFactTitle = contentBoxtf.add_paragraph()
                FunFactTitle.font.name = textFont
                FunFactTitle.font.bold = True
                FunFactTitle.font.size = Pt(textSize)
                FunFactTitle.text = "\nFun Fact:"
                FunFactParagraph = contentBoxtf.add_paragraph()
                FunFactParagraph.font.name = textFont
                FunFactParagraph.font.size = Pt(textSize)
                FunFactParagraph.text = printFact
            #creates the slide presentation if slide option 2 is chosen
            elif(slideOption == 2):
                #creates the slide layout preferences
                slide_layout = prs.slide_layouts[6]
                prs.slide_width = pptx.util.Inches(8)
                prs.slide_height = pptx.util.Inches(11)
                #creates a title page
                if(i == 0):
                    slide = prs.slides.add_slide(slide_layout)
                    titleBox = slide.shapes.add_textbox(pptx.util.Inches(1.5), pptx.util.Inches(2),width=pptx.util.Inches(3), height=pptx.util.Inches(2))
                    titleBoxtf = titleBox.text_frame
                    title = titleBoxtf.add_paragraph()
                    title.text = "Project abcd "
                    title.font.size = Pt(50)
                    title.font.name = titleFont
                slide = prs.slides.add_slide(slide_layout)
                #places the picture to cover the whole slide
                pictureHolder = prs.slides[i+1].shapes
                pictureHolder.add_picture(basename(printImage), pptx.util.Inches(4), pptx.util.Inches(2), width=pptx.util.Inches(4), height=pptx.util.Inches(6))
                
                slideId = slide.shapes.add_textbox(pptx.util.Inches(7), pptx.util.Inches(10), width = pptx.util.Inches(1), height = pptx.util.Inches(1))                       
                slideIdf = slideId.text_frame
                slide_id =slideIdf.add_paragraph()
                slide_id.text =str(value)
                #places the logo on the slide
                logoHolder = slide.shapes.add_picture(basename(printLogo), pptx.util.Inches(7), pptx.util.Inches(10),width=pptx.util.Inches(1), height=pptx.util.Inches(1))
                #places title on the slide
                titleBox = slide.shapes.add_textbox(pptx.util.Inches(2.5), pptx.util.Inches(.5), width=pptx.util.Inches(2),height=pptx.util.Inches(1))
                titleBoxtf = titleBox.text_frame
                title = titleBoxtf.add_paragraph()
                title.text = printName
                title.font.size = Pt(titleSize)
                title.font.name = titleFont
                #creates textbox for description and fun fact
                contentBox = slide.shapes.add_textbox(pptx.util.Inches(1), pptx.util.Inches(.75), width=pptx.util.Inches(3),height=pptx.util.Inches(4))
                contentBoxtf = contentBox.text_frame
                contentBoxtf.word_wrap = True
                descriptionTitle = contentBoxtf.add_paragraph()
                descriptionTitle.font.name = textFont
                descriptionTitle.font.bold = True
                descriptionTitle.font.size = Pt(textSize)
                descriptionTitle.text = "Description: "
                descriptionParagraph = contentBoxtf.add_paragraph()
                descriptionParagraph.font.name = textFont
                descriptionParagraph.font.size = Pt(textSize)
                descriptionParagraph.text = printDescription
                FunFactTitle = contentBoxtf.add_paragraph()
                FunFactTitle.font.bold = True
                FunFactTitle.font.name = textFont
                FunFactTitle.font.size = Pt(textSize)
                FunFactTitle.text = "\nFun Fact:"
                FunFactParagraph = contentBoxtf.add_paragraph()
                FunFactParagraph.font.name = textFont
                FunFactParagraph.font.size = Pt(textSize)
                FunFactParagraph.text = printFact
                
            #creates the slide presentation if slide option 3 is chosen
            elif (slideOption == 3):
                #creates slide preferences
                slide_layout = prs.slide_layouts[6]
                prs.slide_width = pptx.util.Inches(8)
                prs.slide_height = pptx.util.Inches(11)
                slide2 = prs.slides.add_slide(slide_layout)
                #places picture to cover whole slide
                pictureHolder = prs.slides[pictureSlide].shapes
                pictureHolder.add_picture(basename(printImage), pptx.util.Inches(0), pptx.util.Inches(2),width=pptx.util.Inches(4), height=pptx.util.Inches(6))
                
                slideId = slide2.shapes.add_textbox(pptx.util.Inches(7), pptx.util.Inches(10), width = pptx.util.Inches(1), height = pptx.util.Inches(1))                       
                slideIdf = slideId.text_frame
                slide_id =slideIdf.add_paragraph()
                slide_id.text =str(value)
                
                #place logo on the slide
                #places the title
                titleBox = slide2.shapes.add_textbox(pptx.util.Inches(2.5), pptx.util.Inches(.25),width=pptx.util.Inches(2), height=pptx.util.Inches(1))
                titleBoxtf = titleBox.text_frame
                title = titleBoxtf.add_paragraph()
                title.text = printName
                title.font.size = Pt(titleSize)
                title.font.name = titleFont
                #creates textbox for description and fun fact
                contentBox = slide2.shapes.add_textbox(pptx.util.Inches(4), pptx.util.Inches(.75), width=pptx.util.Inches(3),height=pptx.util.Inches(2))
                contentBoxtf = contentBox.text_frame
                contentBoxtf.word_wrap = True
                descriptionTitle = contentBoxtf.add_paragraph()
                descriptionTitle.font.name = textFont
                descriptionTitle.font.bold = True
                descriptionTitle.font.size = Pt(textSize)
                descriptionTitle.text = "Description: "
                descriptionParagraph = contentBoxtf.add_paragraph()
                descriptionParagraph.font.name = textFont
                descriptionParagraph.font.size = Pt(textSize)
                descriptionParagraph.text = printDescription
                FunFactTitle = contentBoxtf.add_paragraph()
                FunFactTitle.font.bold = True
                FunFactTitle.font.name = textFont
                FunFactTitle.font.size = Pt(textSize)
                FunFactTitle.text = "\nFun Fact:"
                FunFactParagraph = contentBoxtf.add_paragraph()
                FunFactParagraph.font.name = textFont
                FunFactParagraph.font.size = Pt(textSize)
                FunFactParagraph.text = printFact
                pictureSlide = pictureSlide + 1
            elif(slideOption == 4):
                #creates the slide layout preferences
                slide_layout = prs.slide_layouts[6]
                prs.slide_width = pptx.util.Inches(8)
                prs.slide_height = pptx.util.Inches(11)
                #creates a title page
                if(i == 0):
                    slide1 = prs.slides.add_slide(slide_layout)
                    titleBox = slide1.shapes.add_textbox(pptx.util.Inches(1.5), pptx.util.Inches(4.5),width=pptx.util.Inches(5), height=pptx.util.Inches(2))
                    titleBoxtf = titleBox.text_frame
                    title = titleBoxtf.add_paragraph()
                    title.text = "Project abcd "
                    title.font.size = Pt(60)
                    title.font.name = titleFont
                slide = prs.slides.add_slide(slide_layout)
                #places the picture to cover the whole slide
                pictureHolder = prs.slides[pictureSlide].shapes
                pictureHolder.add_picture(basename(printImage), pptx.util.Inches(0), pptx.util.Inches(0),width=pptx.util.Inches(8), height=pptx.util.Inches(11))
                #creates next slide
                slide2 = prs.slides.add_slide(slide_layout)
                #places the logo on the slide
                slideId = slide2.shapes.add_textbox(pptx.util.Inches(7), pptx.util.Inches(10), width = pptx.util.Inches(1), height = pptx.util.Inches(1))                       
                slideIdf = slideId.text_frame
                slide_id =slideIdf.add_paragraph()
                slide_id.text =str(value)
                logoHolder = slide2.shapes.add_picture(basename(printLogo), pptx.util.Inches(7), pptx.util.Inches(10),width=pptx.util.Inches(1), height=pptx.util.Inches(1))
                #places title on the slide
                titleBox = slide2.shapes.add_textbox(pptx.util.Inches(2), pptx.util.Inches(.5), width=pptx.util.Inches(4),height=pptx.util.Inches(1))
                titleBoxtf = titleBox.text_frame
                title = titleBoxtf.add_paragraph()
                title.text = printName
                title.font.size = Pt(titleSize)
                title.font.name = titleFont
                #creates textbox for description and fun fact
                contentBox = slide2.shapes.add_textbox(pptx.util.Inches(1), pptx.util.Inches(2), width=pptx.util.Inches(6),height=pptx.util.Inches(7))
                contentBoxtf = contentBox.text_frame
                contentBoxtf.word_wrap = True
                descriptionTitle = contentBoxtf.add_paragraph()
                descriptionTitle.font.name = textFont
                descriptionTitle.font.bold = True
                descriptionTitle.font.size = Pt(textSize)
                descriptionTitle.text = "Description: "
                descriptionParagraph = contentBoxtf.add_paragraph()
                descriptionParagraph.font.name = textFont
                descriptionParagraph.font.size = Pt(textSize)
                descriptionParagraph.text = printDescription
                FunFactTitle = contentBoxtf.add_paragraph()
                FunFactTitle.font.bold = True
                FunFactTitle.font.name = textFont
                FunFactTitle.font.size = Pt(textSize)
                FunFactTitle.text = "\nFun Fact:"
                FunFactParagraph = contentBoxtf.add_paragraph()
                FunFactParagraph.font.name = textFont
                FunFactParagraph.font.size = Pt(textSize)
                FunFactParagraph.text = printFact
                pictureSlide = pictureSlide + 2
            elif(slideOption == 5 ):
                #creates the slides and sets layout preferences
                slide_layout = prs.slide_layouts[6]
                slide = prs.slides.add_slide(slide_layout)
                prs.slide_width = pptx.util.Inches(11)
                prs.slide_height = pptx.util.Inches(8)
                
                slideId = slide.shapes.add_textbox(pptx.util.Inches(7), pptx.util.Inches(10), width = pptx.util.Inches(1), height = pptx.util.Inches(1))                       
                slideIdf = slideId.text_frame
                slide_id =slideIdf.add_paragraph()
                slide_id.text =str(value) 
                #places the logo on the slide
                logoHolder = slide.shapes.add_picture(basename(printLogo), pptx.util.Inches(10), pptx.util.Inches(7),width=pptx.util.Inches(1), height=pptx.util.Inches(1))
                #places the title on the slide
                titleBox = slide.shapes.add_textbox(pptx.util.Inches(3.5), pptx.util.Inches(.25),width=pptx.util.Inches(2), height=pptx.util.Inches(1))
                titleBoxtf = titleBox.text_frame
                title = titleBoxtf.add_paragraph()
                title.text = printName
                title.font.name = titleFont
                title.font.size = Pt(titleSize)
                #places the picture on the slide
                pictureHolder = prs.slides[i].shapes
                pictureHolder.add_picture(basename(printImage), pptx.util.Inches(1),pptx.util.Inches(1),width = pptx.util.Inches(4), height = pptx.util.Inches(5))
                #creates a textbox for the description and fun fact
                contentBox = slide.shapes.add_textbox(pptx.util.Inches(6), pptx.util.Inches(.75), width=pptx.util.Inches(4),height=pptx.util.Inches(1))
                contentBoxtf = contentBox.text_frame
                contentBoxtf.word_wrap = True
                descriptionTitle = contentBoxtf.add_paragraph()
                descriptionTitle.font.name = textFont
                descriptionTitle.font.bold = True
                descriptionTitle.font.size = Pt(textSize)
                descriptionTitle.text = "Description: "
                descriptionParagraph = contentBoxtf.add_paragraph()
                descriptionParagraph.font.name = textFont
                descriptionParagraph.font.size = Pt(textSize)
                descriptionParagraph.text = printDescription
                
                factBox = slide.shapes.add_textbox(pptx.util.Inches(1), pptx.util.Inches(5), width=pptx.util.Inches(3),height=pptx.util.Inches(1))
                factBoxf = factBox.text_frame
                factBoxf.word_wrap = True
                FunFactTitle = factBoxf.add_paragraph()
                FunFactTitle.font.name = textFont
                FunFactTitle.font.bold = True
                FunFactTitle.font.size = Pt(textSize)
                FunFactTitle.text = "\nFun Fact:"
                FunFactParagraph = factBoxf.add_paragraph()
                FunFactParagraph.font.name = textFont
                FunFactParagraph.font.size = Pt(textSize)
                FunFactParagraph.text = printFact

        test = output +".pptx"
        prs.save(test)
        return test
    test = buildPresentation()
    os.startfile(test)
    
    
elif(methods == "Excel"):
    df = pd.read_excel(r"C:\xampp\htdocs\projectabcd_admin\dresses.xlsx")

    
    def buildPresentation(df):
          
     print("Creating powerpoint slides.")
    
     with open(preference, "r") as f:
                slideOption = int(layout)
                textFont = f.readline().split("= ")
                textFont = textFont[1]
                titleFont = f.readline().split("= ")
                titleFont = titleFont[1]
                textSize = f.readline().split("= ")
                textSize = int(textSize[1])
                titleSize = f.readline().split("= ")
                titleSize = int(titleSize[1])
                prs = Presentation()
                PresentationLength = len(all_pages)
                list_values_ = []
                index_value =[]
                if(sort_choice == "3"):
                    for i in range (len(all_pages)):
                        new_val = int(all_pages[i])
                        new_val = new_val -1
                        name_value = df['name'][new_val]
                        id_value = df['id'][new_val]
                        list_values_.append((id_value,name_value))
                    sorted_list = sorted(list_values_, key=lambda x: x[1])
                    list_values_ = sorted_list
                    for i in range (len(list_values_)):
                        index_value.append(list_values_[i][0])
                if(sort_choice == "2"):
                    for i in range (len(all_pages)):
                        new_val = int(all_pages[i])
                        list_values_.append(new_val)
                        
                    list_values_.sort()
                                        
                
     for i in range (PresentationLength):   
            if(sort_choice == "3"):
                value = int(index_value[i]) 
                value = value -1
            elif(sort_choice == "2"):
                value = int(list_values_[i])
                value = value -1
            elif(sort_choice == "1"):
                value = int(all_pages[i])
                value = value -1
            else:
                print("Choose a sort type")
                sys.exit()
                
            if(slideOption == 1 ):
                               
                    slide = prs.slides.add_slide(prs.slide_layouts[6]) 
                    
                    prs.slide_width = pptx.util.Inches(8)
                    prs.slide_height = pptx.util.Inches(11)        
                    
                    slideId = slide.shapes.add_textbox(pptx.util.Inches(7), pptx.util.Inches(10), width = pptx.util.Inches(1), height = pptx.util.Inches(1))                       
                    slideIdf = slideId.text_frame
                    slide_id =slideIdf.add_paragraph()
                    slide_id.text =str(value+1) 
                    contentBox = slide.shapes.add_textbox(pptx.util.Inches(1), pptx.util.Inches(6), width = pptx.util.Inches(6), height = pptx.util.Inches(5))
                    titleBox = slide.shapes.add_textbox(pptx.util.Inches(4), pptx.util.Inches(.5),width=pptx.util.Inches(2), height=pptx.util.Inches(1))
                    
                    titleBoxtf = titleBox.text_frame
                    title = titleBoxtf.add_paragraph()
                    title.font.name = titleFont
                    title.font.size = Pt(titleSize)
                    title.font.name = titleFont
                    title.text =(df['name'][value]) 
                    
                    contentBoxtf = contentBox.text_frame
                    contentBoxtf.word_wrap = True
                    descriptionTitle = contentBoxtf.add_paragraph()
                    descriptionTitle.font.name = textFont
                    descriptionTitle.font.bold = True
                    descriptionTitle.font.size = Pt(textSize)
                    descriptionTitle.text = "Description: "
                    descriptionParagraph = contentBoxtf.add_paragraph()
                    descriptionParagraph.font.name = textFont
                    descriptionParagraph.font.size = Pt(textSize)
                    descriptionParagraph.text = (df['description'][value])
                    FunFactTitle = contentBoxtf.add_paragraph()
                    FunFactTitle.font.bold = True
                    FunFactTitle.font.name = textFont
                    FunFactTitle.font.size = Pt(textSize)
                    FunFactTitle.text = "\nFun Fact:"
                    FunFactParagraph = contentBoxtf.add_paragraph()
                    FunFactParagraph.font.name = textFont
                    FunFactParagraph.font.size = Pt(textSize)
                    FunFactParagraph.text =  str(df['did_you_know'][value])
                    
                    
                    
                    image_url = (df['image_url'][value]) 
                    if image_url:
                        image_path = os.path.basename(image_url)
                        slide.shapes.add_picture(image_path, pptx.util.Inches(2.5),pptx.util.Inches(2),width = pptx.util.Inches(3), height = pptx.util.Inches(4))

            elif(slideOption == 2):     
                    slide = prs.slides.add_slide(prs.slide_layouts[6]) 
                    
                    prs.slide_width = pptx.util.Inches(8)
                    prs.slide_height = pptx.util.Inches(11)        
                    
                    slideId = slide.shapes.add_textbox(pptx.util.Inches(7), pptx.util.Inches(10), width = pptx.util.Inches(1), height = pptx.util.Inches(1))                       
                    slideIdf = slideId.text_frame
                    slide_id =slideIdf.add_paragraph()
                    slide_id.text =str(value+1) 
                    contentBox = slide.shapes.add_textbox(pptx.util.Inches(1), pptx.util.Inches(.75), width=pptx.util.Inches(3),height=pptx.util.Inches(4))
                    titleBox = slide.shapes.add_textbox(pptx.util.Inches(3.5), pptx.util.Inches(.5), width=pptx.util.Inches(2),height=pptx.util.Inches(1))
                    
                    titleBoxtf = titleBox.text_frame
                    title = titleBoxtf.add_paragraph()
                    title.font.name = titleFont
                    title.font.size = Pt(titleSize)
                    title.font.name = titleFont
                    title.text = (df['name'][value]) 
                    
                    contentBoxtf = contentBox.text_frame
                    contentBoxtf.word_wrap = True
                    descriptionTitle = contentBoxtf.add_paragraph()
                    descriptionTitle.font.name = textFont
                    descriptionTitle.font.bold = True
                    descriptionTitle.font.size = Pt(textSize)
                    descriptionTitle.text = "Description: "
                    descriptionParagraph = contentBoxtf.add_paragraph()
                    descriptionParagraph.font.name = textFont
                    descriptionParagraph.font.size = Pt(textSize)
                    descriptionParagraph.text = (df['description'][value])
                    FunFactTitle = contentBoxtf.add_paragraph()
                    FunFactTitle.font.bold = True
                    FunFactTitle.font.name = textFont
                    FunFactTitle.font.size = Pt(textSize)
                    FunFactTitle.text = "\nFun Fact:"
                    FunFactParagraph = contentBoxtf.add_paragraph()
                    FunFactParagraph.font.name = textFont
                    FunFactParagraph.font.size = Pt(textSize)
                    FunFactParagraph.text = str(df['did_you_know'][value])
                    
                    
                    
                    image_url = (df['image_url'][value]) 
                    if image_url:
                        image_path = os.path.basename(image_url)
                        slide.shapes.add_picture(image_path, pptx.util.Inches(4), pptx.util.Inches(2), width=pptx.util.Inches(4), height=pptx.util.Inches(6))
        
            elif(slideOption == 3):
                    slide = prs.slides.add_slide(prs.slide_layouts[6]) 
                    
                    prs.slide_width = pptx.util.Inches(8)
                    prs.slide_height = pptx.util.Inches(11)        
                    
                    slideId = slide.shapes.add_textbox(pptx.util.Inches(7), pptx.util.Inches(10), width = pptx.util.Inches(1), height = pptx.util.Inches(1))                       
                    slideIdf = slideId.text_frame
                    slide_id =slideIdf.add_paragraph()
                    slide_id.text =str(value+1) 
                    contentBox = slide.shapes.add_textbox(pptx.util.Inches(4), pptx.util.Inches(.75), width=pptx.util.Inches(3),height=pptx.util.Inches(2))
                    titleBox = slide.shapes.add_textbox(pptx.util.Inches(3.5), pptx.util.Inches(.25),width=pptx.util.Inches(2), height=pptx.util.Inches(1))
                    
                    titleBoxtf = titleBox.text_frame
                    title = titleBoxtf.add_paragraph()
                    title.font.name = titleFont
                    title.font.size = Pt(titleSize)
                    title.font.name = titleFont
                    title.text = (df['name'][value]) 
                    
                    contentBoxtf = contentBox.text_frame
                    contentBoxtf.word_wrap = True
                    descriptionTitle = contentBoxtf.add_paragraph()
                    descriptionTitle.font.name = textFont
                    descriptionTitle.font.bold = True
                    descriptionTitle.font.size = Pt(textSize)
                    descriptionTitle.text = "Description: "
                    descriptionParagraph = contentBoxtf.add_paragraph()
                    descriptionParagraph.font.name = textFont
                    descriptionParagraph.font.size = Pt(textSize)
                    descriptionParagraph.text = (df['description'][value])
                    FunFactTitle = contentBoxtf.add_paragraph()
                    FunFactTitle.font.bold = True
                    FunFactTitle.font.name = textFont
                    FunFactTitle.font.size = Pt(textSize)
                    FunFactTitle.text = "\nFun Fact:"
                    FunFactParagraph = contentBoxtf.add_paragraph()
                    FunFactParagraph.font.name = textFont
                    FunFactParagraph.font.size = Pt(textSize)
                    FunFactParagraph.text = str(df['did_you_know'][value])
                    
                    
                    
                    image_url = (df['image_url'][value]) 
                    if image_url:
                        image_path = os.path.basename(image_url)
                        slide.shapes.add_picture(image_path, pptx.util.Inches(0), pptx.util.Inches(2),width=pptx.util.Inches(4), height=pptx.util.Inches(6))
            
            elif(slideOption == 4):
                    slide = prs.slides.add_slide(prs.slide_layouts[6]) 
                    slide2 = prs.slides.add_slide(prs.slide_layouts[6])
                    
                    prs.slide_width = pptx.util.Inches(8)
                    prs.slide_height = pptx.util.Inches(11)        
                    
                    slideId = slide.shapes.add_textbox(pptx.util.Inches(7), pptx.util.Inches(10), width = pptx.util.Inches(1), height = pptx.util.Inches(1))                       
                    slideIdf = slideId.text_frame
                    slide_id =slideIdf.add_paragraph()
                    slide_id.text =str(value+1) 
                    contentBox = slide.shapes.add_textbox(pptx.util.Inches(6), pptx.util.Inches(.75), width=pptx.util.Inches(4),height=pptx.util.Inches(1))
                    titleBox = slide.shapes.add_textbox(pptx.util.Inches(3.5), pptx.util.Inches(.25),width=pptx.util.Inches(2), height=pptx.util.Inches(1)) 
                    titleBoxtf = titleBox.text_frame
                    title = titleBoxtf.add_paragraph()
                    title.font.name = titleFont
                    title.font.size = Pt(titleSize)
                    title.font.name = titleFont
                    title.text = (df['name'][value]) 
                    
                    contentBoxtf = contentBox.text_frame
                    contentBoxtf.word_wrap = True
                    descriptionTitle = contentBoxtf.add_paragraph()
                    descriptionTitle.font.name = textFont
                    descriptionTitle.font.bold = True
                    descriptionTitle.font.size = Pt(textSize)
                    descriptionTitle.text = "Description: "
                    descriptionParagraph = contentBoxtf.add_paragraph()
                    descriptionParagraph.font.name = textFont
                    descriptionParagraph.font.size = Pt(textSize)
                    descriptionParagraph.text = (df['description'][value])
                    FunFactTitle = contentBoxtf.add_paragraph()
                    FunFactTitle.font.bold = True
                    FunFactTitle.font.name = textFont
                    FunFactTitle.font.size = Pt(textSize)
                    FunFactTitle.text = "\nFun Fact:"
                    FunFactParagraph = contentBoxtf.add_paragraph()
                    FunFactParagraph.font.name = textFont
                    FunFactParagraph.font.size = Pt(textSize)
                    FunFactParagraph.text = str(df['did_you_know'][value])
                    
                    image_url = df['image_url'][value]
                    if image_url:
                        image_path = os.path.basename(image_url)
                        slide.shapes.add_picture(image_path, pptx.util.Inches(0), pptx.util.Inches(0),width=pptx.util.Inches(8), height=pptx.util.Inches(11))
                            
            if(slideOption == 5 ):
                               
                    slide = prs.slides.add_slide(prs.slide_layouts[6]) 
                    
                    prs.slide_width = pptx.util.Inches(11)
                    prs.slide_height = pptx.util.Inches(8)        
                    
                    slideId = slide.shapes.add_textbox(pptx.util.Inches(7), pptx.util.Inches(10), width = pptx.util.Inches(1), height = pptx.util.Inches(1))                       
                    slideIdf = slideId.text_frame
                    slide_id =slideIdf.add_paragraph()
                    slide_id.text =str(value+1) 
                    contentBox = slide.shapes.add_textbox(pptx.util.Inches(6), pptx.util.Inches(.75), width = pptx.util.Inches(4), height = pptx.util.Inches(1))
                    titleBox = slide.shapes.add_textbox(pptx.util.Inches(3.5), pptx.util.Inches(.25),width=pptx.util.Inches(2), height=pptx.util.Inches(1))
                    
                    titleBoxtf = titleBox.text_frame
                    title = titleBoxtf.add_paragraph()
                    title.font.name = titleFont
                    title.font.size = Pt(titleSize)
                    title.font.name = titleFont
                    title.text =(df['name'][value]) 
                    
                    contentBoxtf = contentBox.text_frame
                    contentBoxtf.word_wrap = True
                    descriptionTitle = contentBoxtf.add_paragraph()
                    descriptionTitle.font.name = textFont
                    descriptionTitle.font.bold = True
                    descriptionTitle.font.size = Pt(textSize)
                    descriptionTitle.text = "Description: "
                    descriptionParagraph = contentBoxtf.add_paragraph()
                    descriptionParagraph.font.name = textFont
                    descriptionParagraph.font.size = Pt(textSize)
                    descriptionParagraph.text = (df['description'][value])
                    factBox = slide.shapes.add_textbox(pptx.util.Inches(1), pptx.util.Inches(5), width=pptx.util.Inches(3),height=pptx.util.Inches(1))
                    factBoxf = factBox.text_frame
                    factBoxf.word_wrap = True
                    FunFactTitle = factBoxf.add_paragraph()
                    FunFactTitle.font.bold = True
                    FunFactTitle.font.name = textFont
                    FunFactTitle.font.size = Pt(textSize)
                    FunFactTitle.text = "\nFun Fact:"
                    FunFactParagraph = factBoxf.add_paragraph()
                    FunFactParagraph.font.name = textFont
                    FunFactParagraph.font.size = Pt(textSize)
                    FunFactParagraph.text =  str(df['did_you_know'][value])
                
                    image_url = (df['image_url'][value]) 
                    if image_url:
                        image_path = os.path.basename(image_url)
                        slide.shapes.add_picture(image_path, pptx.util.Inches(1),pptx.util.Inches(1),width = pptx.util.Inches(4), height = pptx.util.Inches(5))
       
                                    
     test = output +".pptx"
     prs.save(test)
     return test

    test = buildPresentation(df)
    os.startfile(test)

elif(methods == "Database"):
    conn = mysql.connector.connect(user='root', host='localhost', database='abcd_dress-500')
    cursor = conn.cursor()


    sql = "SELECT * FROM dresses"
    cursor.execute(sql)
    result = cursor.fetchall()

    def buildPresentation(data):
        print("Creating PowerPoint slides.")

        with open("preferences.txt", "r") as f:
            slideOption = int(layout)
            textFont = f.readline().split("= ")
            textFont = textFont[1]
            titleFont = f.readline().split("= ")
            titleFont = titleFont[1]
            textSize = f.readline().split("= ")
            textSize = int(textSize[1])
            titleSize = f.readline().split("= ")
            titleSize = int(titleSize[1])

        prs = Presentation()
        PresentationLength = len(all_pages)
       
        list_values_ = []
        index_value =[]
        if(sort_choice == "3"):
            for i in range (len(all_pages)):
                new_val = int(all_pages[i])
                new_val = new_val -1
                name_value = (data[new_val][1])
                id_value = (data[new_val][0])
                list_values_.append((id_value,name_value))
            sorted_list = sorted(list_values_, key=lambda x: x[1])
            list_values_ = sorted_list
            for i in range (len(list_values_)):
                index_value.append(list_values_[i][0])

        if(sort_choice == "2"):
            for i in range (len(all_pages)):
                new_val = int(all_pages[i])
                list_values_.append(new_val)
            list_values_.sort()
                
        for i in range (PresentationLength):   
            if(sort_choice == "3"):
                value = int(index_value[i])
                value = value -1
            elif(sort_choice == "2"):
                value = int(list_values_[i])
                value = value -1
            elif(sort_choice == "1"):
                value = int(all_pages[i])
                value = value -1
            else:
                print("Choose a sort type")
                sys.exit()
                
            if(slideOption == 1 ):
                                   
                        slide = prs.slides.add_slide(prs.slide_layouts[6]) 
                        prs.slide_width = pptx.util.Inches(8)
                        prs.slide_height = pptx.util.Inches(11)        
                        slideId = slide.shapes.add_textbox(pptx.util.Inches(7), pptx.util.Inches(10), width = pptx.util.Inches(1), height = pptx.util.Inches(1))                       
                        slideIdf = slideId.text_frame
                        slide_id =slideIdf.add_paragraph()
                        slide_id.text =str(value+1)
                        contentBox = slide.shapes.add_textbox(pptx.util.Inches(1), pptx.util.Inches(6), width = pptx.util.Inches(6), height = pptx.util.Inches(5))
                        titleBox = slide.shapes.add_textbox(pptx.util.Inches(3.5), pptx.util.Inches(.5),width=pptx.util.Inches(2), height=pptx.util.Inches(1))
                        
                        titleBoxtf = titleBox.text_frame
                        title = titleBoxtf.add_paragraph()
                        title.font.name = titleFont
                        title.font.size = Pt(titleSize)
                        title.font.name = titleFont
                        title.text = (data[value][1]) 
                        
                        contentBoxtf = contentBox.text_frame
                        contentBoxtf.word_wrap = True
                        descriptionTitle = contentBoxtf.add_paragraph()
                        descriptionTitle.font.name = textFont
                        descriptionTitle.font.bold = True
                        descriptionTitle.font.size = Pt(textSize)
                        descriptionTitle.text = "Description: "
                        descriptionParagraph = contentBoxtf.add_paragraph()
                        descriptionParagraph.font.name = textFont
                        descriptionParagraph.font.size = Pt(textSize)
                        descriptionParagraph.text = (data[value][2])
                        FunFactTitle = contentBoxtf.add_paragraph()
                        FunFactTitle.font.bold = True
                        FunFactTitle.font.name = textFont
                        FunFactTitle.font.size = Pt(textSize)
                        FunFactTitle.text = "\nFun Fact:"
                        FunFactParagraph = contentBoxtf.add_paragraph()
                        FunFactParagraph.font.name = textFont
                        FunFactParagraph.font.size = Pt(textSize)
                        FunFactParagraph.text =  (data[value][3])
                        
                        
                        
                        image_url = (data[value][8] )
                        if image_url:
                            image_path = os.path.basename(image_url)
                            slide.shapes.add_picture(image_path, pptx.util.Inches(2.5),pptx.util.Inches(2),width = pptx.util.Inches(3), height = pptx.util.Inches(4))

            elif(slideOption == 2):     
                        slide = prs.slides.add_slide(prs.slide_layouts[6]) 
                        
                        prs.slide_width = pptx.util.Inches(8)
                        prs.slide_height = pptx.util.Inches(11)        
                        slideId = slide.shapes.add_textbox(pptx.util.Inches(7), pptx.util.Inches(10), width = pptx.util.Inches(1), height = pptx.util.Inches(1))                       
                        slideIdf = slideId.text_frame
                        slide_id =slideIdf.add_paragraph()
                        slide_id.text =str(value+1)
                        contentBox = slide.shapes.add_textbox(pptx.util.Inches(1), pptx.util.Inches(.75), width=pptx.util.Inches(3),height=pptx.util.Inches(4))
                        titleBox = slide.shapes.add_textbox(pptx.util.Inches(3.5), pptx.util.Inches(.5),width=pptx.util.Inches(2), height=pptx.util.Inches(1))
                        
                        titleBoxtf = titleBox.text_frame
                        title = titleBoxtf.add_paragraph()
                        title.font.name = titleFont
                        title.font.size = Pt(titleSize)
                        title.font.name = titleFont
                        title.text = (data[value][1]) 
                        
                        contentBoxtf = contentBox.text_frame
                        contentBoxtf.word_wrap = True
                        descriptionTitle = contentBoxtf.add_paragraph()
                        descriptionTitle.font.name = textFont
                        descriptionTitle.font.bold = True
                        descriptionTitle.font.size = Pt(textSize)
                        descriptionTitle.text = "Description: "
                        descriptionParagraph = contentBoxtf.add_paragraph()
                        descriptionParagraph.font.name = textFont
                        descriptionParagraph.font.size = Pt(textSize)
                        descriptionParagraph.text = (data[value][2])
                        FunFactTitle = contentBoxtf.add_paragraph()
                        FunFactTitle.font.bold = True
                        FunFactTitle.font.name = textFont
                        FunFactTitle.font.size = Pt(textSize)
                        FunFactTitle.text = "\nFun Fact:"
                        FunFactParagraph = contentBoxtf.add_paragraph()
                        FunFactParagraph.font.name = textFont
                        FunFactParagraph.font.size = Pt(textSize)
                        FunFactParagraph.text = (data[value][3])
                        
                        
                        
                        image_url = image_url = data[value][8] 
                        if image_url:
                            image_path = os.path.basename(image_url)
                            slide.shapes.add_picture(image_path, pptx.util.Inches(4), pptx.util.Inches(2), width=pptx.util.Inches(4), height=pptx.util.Inches(6))
            
            elif(slideOption == 3):
                        slide = prs.slides.add_slide(prs.slide_layouts[6]) 
                        
                        prs.slide_width = pptx.util.Inches(8)
                        prs.slide_height = pptx.util.Inches(11)        
                        
                        slideId = slide.shapes.add_textbox(pptx.util.Inches(7), pptx.util.Inches(10), width = pptx.util.Inches(1), height = pptx.util.Inches(1))                       
                        slideIdf = slideId.text_frame
                        slide_id =slideIdf.add_paragraph()
                        slide_id.text =str(value+1)
                        contentBox = slide.shapes.add_textbox(pptx.util.Inches(4), pptx.util.Inches(.75), width=pptx.util.Inches(3),height=pptx.util.Inches(2))
                        titleBox = slide.shapes.add_textbox(pptx.util.Inches(3.5), pptx.util.Inches(.5),width=pptx.util.Inches(2), height=pptx.util.Inches(1))
                        
                        titleBoxtf = titleBox.text_frame
                        title = titleBoxtf.add_paragraph()
                        title.font.name = titleFont
                        title.font.size = Pt(titleSize)
                        title.font.name = titleFont
                        title.text = (data[value][1]) 
                        
                        contentBoxtf = contentBox.text_frame
                        contentBoxtf.word_wrap = True
                        descriptionTitle = contentBoxtf.add_paragraph()
                        descriptionTitle.font.name = textFont
                        descriptionTitle.font.bold = True
                        descriptionTitle.font.size = Pt(textSize)
                        descriptionTitle.text = "Description: "
                        descriptionParagraph = contentBoxtf.add_paragraph()
                        descriptionParagraph.font.name = textFont
                        descriptionParagraph.font.size = Pt(textSize)
                        descriptionParagraph.text = (data[i][2])
                        FunFactTitle = contentBoxtf.add_paragraph()
                        FunFactTitle.font.bold = True
                        FunFactTitle.font.name = textFont
                        FunFactTitle.font.size = Pt(textSize)
                        FunFactTitle.text = "\nFun Fact:"
                        FunFactParagraph = contentBoxtf.add_paragraph()
                        FunFactParagraph.font.name = textFont
                        FunFactParagraph.font.size = Pt(textSize)
                        FunFactParagraph.text = (data[value][3])
                        
                        
                        
                        image_url = data[value][8]  
                        if image_url:
                            image_path = os.path.basename(image_url)
                            slide.shapes.add_picture(image_path, pptx.util.Inches(0), pptx.util.Inches(2),width=pptx.util.Inches(4), height=pptx.util.Inches(6))
           
            elif(slideOption == 4):
                    slide = prs.slides.add_slide(prs.slide_layouts[6]) 
                    slide2 = prs.slides.add_slide(prs.slide_layouts[6])
                    
                    prs.slide_width = pptx.util.Inches(8)
                    prs.slide_height = pptx.util.Inches(11)        
                    
                    slideId = slide.shapes.add_textbox(pptx.util.Inches(7), pptx.util.Inches(10), width = pptx.util.Inches(1), height = pptx.util.Inches(1))                       
                    slideIdf = slideId.text_frame
                    slide_id =slideIdf.add_paragraph()
                    slide_id.text =str(value+1)                    
                    contentBox = slide2.shapes.add_textbox(pptx.util.Inches(1), pptx.util.Inches(2), width=pptx.util.Inches(6),height=pptx.util.Inches(7))
                    titleBox = slide.shapes.add_textbox(pptx.util.Inches(3.5), pptx.util.Inches(.5),width=pptx.util.Inches(2), height=pptx.util.Inches(1))
                    
                    titleBoxtf = titleBox.text_frame
                    title = titleBoxtf.add_paragraph()
                    title.font.name = titleFont
                    title.font.size = Pt(titleSize)
                    title.font.name = titleFont
                    title.text = (data[value][1])
                    
                    contentBoxtf = contentBox.text_frame
                    contentBoxtf.word_wrap = True
                    descriptionTitle = contentBoxtf.add_paragraph()
                    descriptionTitle.font.name = textFont
                    descriptionTitle.font.bold = True
                    descriptionTitle.font.size = Pt(textSize)
                    descriptionTitle.text = "Description: "
                    descriptionParagraph = contentBoxtf.add_paragraph()
                    descriptionParagraph.font.name = textFont
                    descriptionParagraph.font.size = Pt(textSize)
                    descriptionParagraph.text = (data[i][2])
                    FunFactTitle = contentBoxtf.add_paragraph()
                    FunFactTitle.font.bold = True
                    FunFactTitle.font.name = textFont
                    FunFactTitle.font.size = Pt(textSize)
                    FunFactTitle.text = "\nFun Fact:"
                    FunFactParagraph = contentBoxtf.add_paragraph()
                    FunFactParagraph.font.name = textFont
                    FunFactParagraph.font.size = Pt(textSize)
                    FunFactParagraph.text = (data[value][3])
                    
                    image_url = data[value][8]  
                    if image_url:
                        image_path = os.path.basename(image_url)
                        slide.shapes.add_picture(image_path, pptx.util.Inches(0), pptx.util.Inches(0),width=pptx.util.Inches(8), height=pptx.util.Inches(11))

            if(slideOption == 5 ):
                               
                    slide = prs.slides.add_slide(prs.slide_layouts[6]) 
                    
                    prs.slide_width = pptx.util.Inches(11)
                    prs.slide_height = pptx.util.Inches(8)        
                    
                    slideId = slide.shapes.add_textbox(pptx.util.Inches(7), pptx.util.Inches(10), width = pptx.util.Inches(1), height = pptx.util.Inches(1))                       
                    slideIdf = slideId.text_frame
                    slide_id =slideIdf.add_paragraph()
                    slide_id.text =str(value+1) 
                    
                    contentBox = slide.shapes.add_textbox(pptx.util.Inches(6), pptx.util.Inches(.75), width = pptx.util.Inches(4), height = pptx.util.Inches(1))
                    titleBox = slide.shapes.add_textbox(pptx.util.Inches(4), pptx.util.Inches(.25),width=pptx.util.Inches(2), height=pptx.util.Inches(1))
                    titleBoxtf = titleBox.text_frame
                    title = titleBoxtf.add_paragraph()
                    title.font.name = titleFont
                    title.font.size = Pt(titleSize)
                    title.font.name = titleFont
                    title.text =(data[value][1]) 
                    
                    contentBoxtf = contentBox.text_frame
                    contentBoxtf.word_wrap = True
                    descriptionTitle = contentBoxtf.add_paragraph()
                    descriptionTitle.font.name = textFont
                    descriptionTitle.font.bold = True
                    descriptionTitle.font.size = Pt(textSize)
                    descriptionTitle.text = "Description: "
                    descriptionParagraph = contentBoxtf.add_paragraph()
                    descriptionParagraph.font.name = textFont
                    descriptionParagraph.font.size = Pt(textSize)
                    descriptionParagraph.text = (data[value][2])
                    factBox = slide.shapes.add_textbox(pptx.util.Inches(1), pptx.util.Inches(5), width=pptx.util.Inches(3),height=pptx.util.Inches(1))
                    factBoxf = factBox.text_frame
                    factBoxf.word_wrap = True
                    FunFactTitle = factBoxf.add_paragraph()
                    FunFactTitle.font.bold = True
                    FunFactTitle.font.name = textFont
                    FunFactTitle.font.size = Pt(textSize)
                    FunFactTitle.text = "\nFun Fact:"
                    FunFactParagraph = factBoxf.add_paragraph()
                    FunFactParagraph.font.name = textFont
                    FunFactParagraph.font.size = Pt(textSize)
                    FunFactParagraph.text =  (data[value][3])
                
                    image_url = (data[value][8])
                    if image_url:
                        image_path = os.path.basename(image_url)
                        slide.shapes.add_picture(image_path, pptx.util.Inches(1),pptx.util.Inches(1),width = pptx.util.Inches(4), height = pptx.util.Inches(5))
       
                                                                               
            
        test = output +".pptx"
        prs.save(test)
        return test
    test = buildPresentation(result)
    os.startfile(test)