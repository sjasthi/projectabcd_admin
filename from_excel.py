from pptx import Presentation
from pptx.util import Pt
import pandas as pd
from sqlalchemy import create_engine
import os
import pptx.util
from collections.abc import Container

df = pd.read_excel(r"C:\xampp\htdocs\projectabcd_admin\dresses.xlsx")
#print(df.columns)
   
def buildPresentation(df):
   start = int(input("Enter the slide to start with: "))
   end = int(input("Enter the slide to end with: "))
   print("Creating powerpoint slides.")
   
   with open("preferences.txt", "r") as f:
        slideOption = f.readline().split("= ")
        slideOption = int(slideOption[1])
        textFont = f.readline().split("= ")
        textFont = textFont[1]
        titleFont = f.readline().split("= ")
        titleFont = titleFont[1]
        textSize = f.readline().split("= ")
        textSize = int(textSize[1])
        titleSize = f.readline().split("= ")
        titleSize = int(titleSize[1])
        prs = Presentation()
        presentationLength = end - start + 1
        start = start -1

   for i in range(0,presentationLength): 
  
          slide = prs.slides.add_slide(prs.slide_layouts[6]) 
          slide2 = prs.slides.add_slide(prs.slide_layouts[6])
        
          prs.slide_width = pptx.util.Inches(8)
          prs.slide_height = pptx.util.Inches(11)        
        
          contentBox = slide2.shapes.add_textbox(pptx.util.Inches(1), pptx.util.Inches(2), width=pptx.util.Inches(6),height=pptx.util.Inches(7))
          titleBox = slide2.shapes.add_textbox(pptx.util.Inches(4), pptx.util.Inches(.5),width=pptx.util.Inches(4), height=pptx.util.Inches(1))
        
          titleBoxtf = titleBox.text_frame
          title = titleBoxtf.add_paragraph()
          title.font.name = titleFont
          title.font.size = Pt(titleSize)
          title.font.name = titleFont
          title.text = str(df['name'][start + i]) 
        
          contentBoxtf = contentBox.text_frame
          contentBoxtf.word_wrap = True
          descriptionTitle = contentBoxtf.add_paragraph()
          descriptionTitle.font.name = textFont
          descriptionTitle.font.bold = True
          descriptionTitle.font.size = Pt(textSize)
          descriptionTitle.text = "Description: "
          descriptionParagraph = contentBoxtf.add_paragraph()
          descriptionParagraph.font.name = textFont
          descriptionParagraph.font.size = Pt(textSize)
          descriptionParagraph.text = str(df['description'][start +i])
          FunFactTitle = contentBoxtf.add_paragraph()
          FunFactTitle.font.bold = True
          FunFactTitle.font.name = textFont
          FunFactTitle.font.size = Pt(textSize)
          FunFactTitle.text = "\nFun Fact:"
          FunFactParagraph = contentBoxtf.add_paragraph()
          FunFactParagraph.font.name = textFont
          FunFactParagraph.font.size = Pt(textSize)
          FunFactParagraph.text = str(df['did_you_know'][start +i])
        
        
        
          image_url = str(df['image_url'][start +i]) 
          if image_url:
               image_path = os.path.basename(image_url)
               slide.shapes.add_picture(image_path, pptx.util.Inches(0), pptx.util.Inches(0),width=pptx.util.Inches(8), height=pptx.util.Inches(11))
        
   test = "test_excel.pptx"
   prs.save(test)
   return test

test = buildPresentation(df)
os.startfile(test)
