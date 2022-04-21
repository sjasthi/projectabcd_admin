import os
import shutil
from os.path import basename

import pptx.util
from pptx import Presentation
from pptx.util import Pt
from bs4 import BeautifulSoup
import requests





def buildPresentation():
   with open("preferences.txt", "r") as f:
        slideOption = f.readline().split("= ")
        slideOption = int(slideOption[1])
        textFont = f.readline().split("= ")
        textFont = textFont[1]
        titleFont = f.readline().split("= ")
        titleFont = titleFont[1]
        textSize = f.readline().split("= ")
        textSize = int(textSize[1])
        titleSize = f.readline().split("= ")
        titleSize = int(titleSize[1])
   start = int(input("Enter the slide number you'd like to start with: "))
   end = int(input("Enter the slide number you'd like to end with: "))
   prs = Presentation()
   presentationLength = end - start + 1
   pictureSlide = 0
   if(slideOption == 2):
        pictureSlide = 1
   for i in range (0, presentationLength):
        URL = "https://projectabcd.com/display_the_dress.php?id=" + str(start + i)
        page = requests.get(URL, headers={"User-Agent": "html"})
        soup = BeautifulSoup(page.content, "html.parser")
        pageInfo = soup.find("div", class_="container")
        name = pageInfo.find("h2", class_="head")
        printName = name.text
        image = pageInfo.find("image")
        printImage = image.attrs["src"]
        pictureURL = "http://projectabcd.com/" + printImage
        r = requests.get(pictureURL, headers={"User-Agent": "html"}, stream=True)
        if r.status_code == 200:
             with open(basename(printImage), "wb") as f:
                  r.raw.decode_content = True
                  shutil.copyfileobj(r.raw, f)
        element = pageInfo.p
        description = pageInfo.find("p", class_="words")
        printDescription = description.text
        fact = description.find_next_sibling("p")
        printFact = fact.text
        printWords = printDescription + printFact

        if(slideOption == 1 ):
             slide_layout = prs.slide_layouts[6]
             slide = prs.slides.add_slide(slide_layout)
             prs.slide_width = pptx.util.Inches(8)
             prs.slide_height = pptx.util.Inches(11)
             titleBox = slide.shapes.add_textbox(pptx.util.Inches(2.5), pptx.util.Inches(.5), width = pptx.util.Inches(3), height = pptx.util.Inches(1))
             titleBoxtf = titleBox.text_frame
             title = titleBoxtf.add_paragraph()
             title.text = printName
             title.font.name = titleFont
             title.font.size = Pt(titleSize)
             pictureHolder = prs.slides[i].shapes
             pictureHolder.add_picture(basename(printImage), pptx.util.Inches(2.5),pptx.util.Inches(2),width = pptx.util.Inches(3), height = pptx.util.Inches(4))
             contentBox = slide.shapes.add_textbox(pptx.util.Inches(1), pptx.util.Inches(6), width = pptx.util.Inches(6), height = pptx.util.Inches(5))
             contentBoxtf = contentBox.text_frame
             contentBoxtf.word_wrap = True
             descriptionTitle = contentBoxtf.add_paragraph()
             descriptionTitle.font.name = textFont
             descriptionTitle.font.bold = True
             descriptionTitle.font.size = Pt(textSize)
             descriptionTitle.text = "Description: "
             descriptionParagraph = contentBoxtf.add_paragraph()
             descriptionParagraph.font.name = textFont
             descriptionParagraph.font.size = Pt(textSize)
             descriptionParagraph.text = printDescription
             FunFactTitle = contentBoxtf.add_paragraph()
             FunFactTitle.font.name = textFont
             FunFactTitle.font.bold = True
             FunFactTitle.font.size = Pt(textSize)
             FunFactTitle.text = "\nFun Fact:"
             FunFactParagraph = contentBoxtf.add_paragraph()
             FunFactParagraph.font.name = textFont
             FunFactParagraph.font.size = Pt(textSize)
             FunFactParagraph.text = printFact
        elif(slideOption == 2):
             slide_layout = prs.slide_layouts[6]
             prs.slide_width = pptx.util.Inches(8)
             prs.slide_height = pptx.util.Inches(11)
             if(i == 0):
                  slide1 = prs.slides.add_slide(slide_layout)
                  titleBox = slide1.shapes.add_textbox(pptx.util.Inches(1.5), pptx.util.Inches(4.5),width=pptx.util.Inches(5), height=pptx.util.Inches(2))
                  titleBoxtf = titleBox.text_frame
                  title = titleBoxtf.add_paragraph()
                  title.text = "Project abcd"
                  title.font.size = Pt(60)
                  title.font.name = titleFont
             slide = prs.slides.add_slide(slide_layout)
             pictureHolder = prs.slides[pictureSlide].shapes
             pictureHolder.add_picture(basename(printImage), pptx.util.Inches(0), pptx.util.Inches(0),width=pptx.util.Inches(8), height=pptx.util.Inches(11))
             slide2 = prs.slides.add_slide(slide_layout)
             titleBox = slide2.shapes.add_textbox(pptx.util.Inches(2), pptx.util.Inches(.5), width=pptx.util.Inches(4),height=pptx.util.Inches(1))
             titleBoxtf = titleBox.text_frame
             title = titleBoxtf.add_paragraph()
             title.text = printName
             title.font.size = Pt(titleSize)
             title.font.name = titleFont
             contentBox = slide2.shapes.add_textbox(pptx.util.Inches(1), pptx.util.Inches(2), width=pptx.util.Inches(6),height=pptx.util.Inches(7))
             contentBoxtf = contentBox.text_frame
             contentBoxtf.word_wrap = True
             descriptionTitle = contentBoxtf.add_paragraph()
             descriptionTitle.font.name = textFont
             descriptionTitle.font.bold = True
             descriptionTitle.font.size = Pt(textSize)
             descriptionTitle.text = "Description: "
             descriptionParagraph = contentBoxtf.add_paragraph()
             descriptionParagraph.font.name = textFont
             descriptionParagraph.font.size = Pt(textSize)
             descriptionParagraph.text = printDescription
             FunFactTitle = contentBoxtf.add_paragraph()
             FunFactTitle.font.bold = True
             FunFactTitle.font.name = textFont
             FunFactTitle.font.size = Pt(textSize)
             FunFactTitle.text = "\nFun Fact:"
             FunFactParagraph = contentBoxtf.add_paragraph()
             FunFactParagraph.font.name = textFont
             FunFactParagraph.font.size = Pt(textSize)
             FunFactParagraph.text = printFact
             pictureSlide = pictureSlide + 2
        elif (slideOption == 3):
             slide_layout = prs.slide_layouts[6]
             prs.slide_width = pptx.util.Inches(8)
             prs.slide_height = pptx.util.Inches(11)
             slide2 = prs.slides.add_slide(slide_layout)
             titleBox3 = slide2.shapes.add_textbox(pptx.util.Inches(1.5), pptx.util.Inches(4.5),width=pptx.util.Inches(5), height=pptx.util.Inches(2))
             titleBox3tf = titleBox3.text_frame
             title3 = titleBox3tf.add_paragraph()
             title3.text = printName
             title3.font.name = titleFont
             title3.font.size = Pt(titleSize)
             pictureHolder = prs.slides[pictureSlide].shapes
             pictureHolder.add_picture(basename(printImage), pptx.util.Inches(0), pptx.util.Inches(0),width=pptx.util.Inches(8), height=pptx.util.Inches(11))
             slide3 = prs.slides.add_slide(slide_layout)
             titleBox = slide3.shapes.add_textbox(pptx.util.Inches(4), pptx.util.Inches(.5),width=pptx.util.Inches(4), height=pptx.util.Inches(1))
             titleBoxtf = titleBox.text_frame
             title = titleBoxtf.add_paragraph()
             title.text = printName
             title.font.size = Pt(titleSize)
             title.font.name = titleFont
             contentBox = slide3.shapes.add_textbox(pptx.util.Inches(1), pptx.util.Inches(2), width=pptx.util.Inches(6),height=pptx.util.Inches(7))
             contentBoxtf = contentBox.text_frame
             contentBoxtf.word_wrap = True
             descriptionTitle = contentBoxtf.add_paragraph()
             descriptionTitle.font.name = textFont
             descriptionTitle.font.bold = True
             descriptionTitle.font.size = Pt(textSize)
             descriptionTitle.text = "Description: "
             descriptionParagraph = contentBoxtf.add_paragraph()
             descriptionParagraph.font.name = textFont
             descriptionParagraph.font.size = Pt(textSize)
             descriptionParagraph.text = printDescription
             FunFactTitle = contentBoxtf.add_paragraph()
             FunFactTitle.font.bold = True
             FunFactTitle.font.name = textFont
             FunFactTitle.font.size = Pt(textSize)
             FunFactTitle.text = "\nFun Fact:"
             FunFactParagraph = contentBoxtf.add_paragraph()
             FunFactParagraph.font.name = textFont
             FunFactParagraph.font.size = Pt(textSize)
             FunFactParagraph.text = printFact
             pictureSlide = pictureSlide + 2
   test = "test.pptx"
   prs.save(test)
   return test

test = buildPresentation()
os.startfile(test)

