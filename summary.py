import json
import os
import io

import matplotlib.pyplot as plt
from bs4 import BeautifulSoup
import requests
import pptx.util
from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Pt
from sqlalchemy import null
import matplotlib.pyplot as plt


def _init_(self):
    self = dict()
def add(self, key, value):
    self[key] = value
categoryDict = {}
typeDict = {}
keywordDict = {}
stateDict = {}
statusDict = {}
descriptionDict = {}

URL = "https://projectabcd.com/list_dresses.php"
page = requests.get(URL, headers={"User-Agent": "html"})
soup = BeautifulSoup(page.content, "html.parser")
section = soup.find("tr")
while(section.find_next_sibling("tr")!= null):
    next = section.find_next("tr")
    if(next == None):
        break
    id = next.find("td")
    name = id.find_next_sibling("td")
    description = name.find_next_sibling("td")
    if(len(description.text) > 0 and len(description.text) <= 100):
        descriptionType = "0 - 100"
    elif(len(description.text) > 100 and len(description.text) <= 200):
        descriptionType = "101 - 200"
    elif(len(description.text) > 200 and len(description.text) <= 300):
        descriptionType = "201 - 300"
    elif(len(description.text) > 300 and len(description.text) <= 400):
        descriptionType = "301 - 400"
    elif(len(description.text) > 400 and len(description.text) <= 500):
        descriptionType = "401 - 500"
    elif(len(description.text) > 500 and len(description.text) <= 600):
        descriptionType = "501 - 600"
    elif(len(description.text) > 600 and len(description.text) <= 700):
        descriptionType = "601 - 700"
    elif(len(description.text) > 700 and len(description.text) <= 800):
        descriptionType = "701 - 800"
    elif(len(description.text) > 800 and len(description.text) <= 900):
        descriptionType = "801 - 900"
    elif(len(description.text) > 900 and len(description.text) <= 1000):
        descriptionType = "900 - 100"
    if (descriptionType in descriptionDict):
        descriptionDict[descriptionType] += 1
    else:
        descriptionDict[descriptionType] = 1
    fact = description.find_next_sibling("td")
    category = fact.find_next_sibling("td")
    if(category is not None and category != ""):
        categories = category.text.split(",")
    for cat in categories:
        if(cat.lower() in categoryDict):
            categoryDict[cat.lower()] += 1
        else:
            categoryDict[cat.lower()] = 1
    type = category.find_next_sibling("td")
    if (type.text.lower() in typeDict):
        typeDict[type.text.lower()] += 1
    else:
        typeDict[type.text.lower()] = 1
    stateName = type.find_next_sibling("td")
    if (stateName.text.lower() in stateDict):
        stateDict[stateName.text.lower()] += 1
    else:
        stateDict[stateName.text.lower()] = 1
    keywords = stateName.find_next_sibling("td")
    if (keywords is not None and keywords != ""):
        keywordsList = keywords.text.split(",")
    for key in keywordsList:
        if (key.lower() in keywordDict):
            keywordDict[key.lower()] += 1
        else:
            keywordDict[key.lower()] = 1
    status = keywords.find_next_sibling("td")
    if (status.text.lower() in statusDict):
        statusDict[status.text.lower()] += 1
    else:
        statusDict[status.text.lower()] = 1
    section = next


prs = Presentation()
slide_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(slide_layout)
prs.slide_width = pptx.util.Inches(8)
prs.slide_height = pptx.util.Inches(11)
contentBox = slide.shapes.add_textbox(pptx.util.Inches(1), pptx.util.Inches(2), width = pptx.util.Inches(6), height = pptx.util.Inches(5))
contentBoxtf = contentBox.text_frame
contentBoxtf.word_wrap = True
categoryParagraph = contentBoxtf.add_paragraph()
categoryParagraph.font.size = Pt(16)
x, y, cx, cy = pptx.util.Inches(2), pptx.util.Inches(2), pptx.util.Inches(4), pptx.util.Inches(1.5)
shape = slide.shapes.add_table(len(categoryDict.keys()), 2, x, y, cx, cy)
table = shape.table
count = 0
for i in categoryDict:
    if(count == 0):
        cell0 = table.cell(0,0)
        cell0.text = "Category"
        cell1 = table.cell(0,1)
        cell1.text = "Count"
        count += 1
    else:
        cell = table.cell(count,0)
        cell.text = json.dumps(i).strip('"')
        cell2 = table.cell(count,1)
        cell2.text = json.dumps(categoryDict[i])
        count += 1

slidecatGraph = prs.slides.add_slide(slide_layout)
prs.slide_width = pptx.util.Inches(8)
prs.slide_height = pptx.util.Inches(11)
keys = categoryDict.keys()
values = categoryDict.values()
plt.bar(keys, values)
image_stream = io.BytesIO()
plt.savefig(image_stream)
x, y, cx, cy = pptx.util.Inches(2), pptx.util.Inches(2), pptx.util.Inches(6), pptx.util.Inches(4.5)
pic = slidecatGraph.shapes.add_picture(image_stream, x, y, cx, cy)

slide2 = prs.slides.add_slide(slide_layout)
prs.slide_width = pptx.util.Inches(8)
prs.slide_height = pptx.util.Inches(11)
contentBox = slide2.shapes.add_textbox(pptx.util.Inches(1), pptx.util.Inches(2), width = pptx.util.Inches(6), height = pptx.util.Inches(5))
contentBoxtf = contentBox.text_frame
contentBoxtf.word_wrap = True
typeParagraph = contentBoxtf.add_paragraph()
typeParagraph.font.size = Pt(16)
x, y, cx, cy = pptx.util.Inches(2), pptx.util.Inches(2), pptx.util.Inches(4), pptx.util.Inches(1.5)
shape = slide2.shapes.add_table(len(typeDict.keys()), 2, x, y, cx, cy)
table = shape.table
count = 0
for i in typeDict:
    if (count == 0):
        cell0 = table.cell(0, 0)
        cell0.text = "Type"
        cell1 = table.cell(0, 1)
        cell1.text = "Count"
        count += 1
    else:
        cell = table.cell(count,0)
        cell.text = json.dumps(i).strip('"')
        cell2 = table.cell(count,1)
        cell2.text = json.dumps(typeDict[i])
        count += 1

slidetypeGraph = prs.slides.add_slide(slide_layout)
prs.slide_width = pptx.util.Inches(8)
prs.slide_height = pptx.util.Inches(11)
typeKeys = typeDict.keys()
typeValues = typeDict.values()
plt.bar(typeKeys, typeValues)
image_stream = io.BytesIO()
plt.savefig(image_stream)
x, y, cx, cy = pptx.util.Inches(2), pptx.util.Inches(2), pptx.util.Inches(6), pptx.util.Inches(4.5)
pic = slidetypeGraph.shapes.add_picture(image_stream, x, y, cx, cy)

slide3 = prs.slides.add_slide(slide_layout)
prs.slide_width = pptx.util.Inches(8)
prs.slide_height = pptx.util.Inches(11)
contentBox = slide3.shapes.add_textbox(pptx.util.Inches(1), pptx.util.Inches(2), width = pptx.util.Inches(6), height = pptx.util.Inches(5))
contentBoxtf = contentBox.text_frame
contentBoxtf.word_wrap = True
keywordsParagraph = contentBoxtf.add_paragraph()
keywordsParagraph.font.size = Pt(16)
x, y, cx, cy = pptx.util.Inches(2), pptx.util.Inches(2), pptx.util.Inches(4), pptx.util.Inches(1.5)
shape = slide3.shapes.add_table(len(keywordDict.keys()), 2, x, y, cx, cy)
table = shape.table
count = 0
for i in keywordDict:
    if (count == 0):
        cell0 = table.cell(0, 0)
        cell0.text = "Keywords"
        cell1 = table.cell(0, 1)
        cell1.text = "Count"
        count += 1
    else:
        cell = table.cell(count,0)
        cell.text = json.dumps(i).strip('"')
        cell2 = table.cell(count,1)
        cell2.text = json.dumps(keywordDict[i])
        count += 1

slidekeywordGraph = prs.slides.add_slide(slide_layout)
prs.slide_width = pptx.util.Inches(8)
prs.slide_height = pptx.util.Inches(11)
keywordKeys = keywordDict.keys()
keywordValues = keywordDict.values()
plt.bar(keywordKeys, keywordValues)
image_stream = io.BytesIO()
plt.savefig(image_stream)
x, y, cx, cy = pptx.util.Inches(2), pptx.util.Inches(2), pptx.util.Inches(6), pptx.util.Inches(4.5)
pic = slidekeywordGraph.shapes.add_picture(image_stream, x, y, cx, cy)

slide4 = prs.slides.add_slide(slide_layout)
prs.slide_width = pptx.util.Inches(8)
prs.slide_height = pptx.util.Inches(11)
contentBox = slide4.shapes.add_textbox(pptx.util.Inches(1), pptx.util.Inches(2), width = pptx.util.Inches(6), height = pptx.util.Inches(5))
contentBoxtf = contentBox.text_frame
contentBoxtf.word_wrap = True
stateNameParagraph = contentBoxtf.add_paragraph()
stateNameParagraph.font.size = Pt(16)
x, y, cx, cy = pptx.util.Inches(2), pptx.util.Inches(2), pptx.util.Inches(4), pptx.util.Inches(1.5)
shape = slide4.shapes.add_table(len(stateDict.keys()), 2, x, y, cx, cy)
table = shape.table
count = 0
for i in stateDict:
    if (count == 0):
        cell0 = table.cell(0, 0)
        cell0.text = "State Name"
        cell1 = table.cell(0, 1)
        cell1.text = "Count"
        count += 1
    else:
        cell = table.cell(count,0)
        cell.text = json.dumps(i).strip('"')
        cell2 = table.cell(count,1)
        cell2.text = json.dumps(stateDict[i])
        count += 1

slidestateGraph = prs.slides.add_slide(slide_layout)
prs.slide_width = pptx.util.Inches(8)
prs.slide_height = pptx.util.Inches(11)
stateKeys = stateDict.keys()
stateValues = stateDict.values()
plt.bar(stateKeys, stateValues)
image_stream = io.BytesIO()
plt.savefig(image_stream)
x, y, cx, cy = pptx.util.Inches(2), pptx.util.Inches(2), pptx.util.Inches(6), pptx.util.Inches(4.5)
pic = slidestateGraph.shapes.add_picture(image_stream, x, y, cx, cy)

slide5 = prs.slides.add_slide(slide_layout)
prs.slide_width = pptx.util.Inches(8)
prs.slide_height = pptx.util.Inches(11)
contentBox = slide5.shapes.add_textbox(pptx.util.Inches(1), pptx.util.Inches(2), width = pptx.util.Inches(6), height = pptx.util.Inches(5))
contentBoxtf = contentBox.text_frame
contentBoxtf.word_wrap = True
statusParagraph = contentBoxtf.add_paragraph()
statusParagraph.font.size = Pt(16)
x, y, cx, cy = pptx.util.Inches(2), pptx.util.Inches(2), pptx.util.Inches(4), pptx.util.Inches(1.5)
shape = slide5.shapes.add_table(len(statusDict.keys()), 2, x, y, cx, cy)
table = shape.table
count = 0
for i in statusDict:
    if (count == 0):
        cell0 = table.cell(0, 0)
        cell0.text = "Status"
        cell1 = table.cell(0, 1)
        cell1.text = "Count"
        count += 1
    else:
        cell = table.cell(count,0)
        cell.text = json.dumps(i).strip('"')
        cell2 = table.cell(count,1)
        cell2.text = json.dumps(statusDict[i])
        count += 1

slidestatusGraph = prs.slides.add_slide(slide_layout)
prs.slide_width = pptx.util.Inches(8)
prs.slide_height = pptx.util.Inches(11)
statusKeys = statusDict.keys()
statusValues = statusDict.values()
plt.bar(statusKeys, statusValues)
image_stream = io.BytesIO()
plt.savefig(image_stream)
x, y, cx, cy = pptx.util.Inches(2), pptx.util.Inches(2), pptx.util.Inches(6), pptx.util.Inches(4.5)
pic = slidestatusGraph.shapes.add_picture(image_stream, x, y, cx, cy)

slide6 = prs.slides.add_slide(slide_layout)
prs.slide_width = pptx.util.Inches(8)
prs.slide_height = pptx.util.Inches(11)
contentBox = slide6.shapes.add_textbox(pptx.util.Inches(1), pptx.util.Inches(2), width = pptx.util.Inches(6), height = pptx.util.Inches(5))
contentBoxtf = contentBox.text_frame
contentBoxtf.word_wrap = True
descriptionParagraph = contentBoxtf.add_paragraph()
descriptionParagraph.font.size = Pt(16)
x, y, cx, cy = pptx.util.Inches(2), pptx.util.Inches(2), pptx.util.Inches(4), pptx.util.Inches(1.5)
shape = slide6.shapes.add_table(len(descriptionDict.keys()), 2, x, y, cx, cy)
table = shape.table
for i in descriptionDict:
    cell0 = table.cell(0, 0)
    cell0.text = "Descriptione Character Count"
    cell1 = table.cell(0, 1)
    cell1.text = "Count"

    cell = table.cell(list(descriptionDict).index(i), 0)
    cell.text = json.dumps(i).strip('"')
    cell2 = table.cell(list(descriptionDict).index(i), 1)
    cell2.text = json.dumps(descriptionDict[i])

slidedescriptionGraph = prs.slides.add_slide(slide_layout)
prs.slide_width = pptx.util.Inches(8)
prs.slide_height = pptx.util.Inches(11)
desKeys = descriptionDict.keys()
desValues = descriptionDict.values()
plt.bar(desKeys, desValues)
image_stream = io.BytesIO()
plt.savefig(image_stream)
x, y, cx, cy = pptx.util.Inches(2), pptx.util.Inches(2), pptx.util.Inches(6), pptx.util.Inches(4.5)
pic = slidedescriptionGraph.shapes.add_picture(image_stream, x, y, cx, cy)

summary = "summary.pptx"
prs.save(summary)
os.startfile(summary)