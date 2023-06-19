import os
import shutil
from os.path import basename
from collections.abc import Container
import pptx.util
from pptx import Presentation
from pptx.util import Pt
from bs4 import BeautifulSoup
import requests





def buildPresentation():
   start = int(input("Enter the slide to start with: "))
   end = int(input("Enter the slide to end with: "))
   print("Creating powerpoint slides.")
   #reads preferences text file to get slide orientation option, font, and text size options
   with open("preferences.txt", "r") as f:
        slideOption = f.readline().split("= ")
        slideOption = int(slideOption[1])
        textFont = f.readline().split("= ")
        textFont = textFont[1]
        titleFont = f.readline().split("= ")
        titleFont = titleFont[1]
        textSize = f.readline().split("= ")
        textSize = int(textSize[1])
        titleSize = f.readline().split("= ")
        titleSize = int(titleSize[1])
   prs = Presentation()
   presentationLength = end - start + 1
   pictureSlide = 0
   if(slideOption == 2):
        pictureSlide = 1
   #web scrapes the URL to get all needed information from the page
   for i in range (0, presentationLength):
        URL = "https://projectabcd.com/display_the_dress.php?id=" + str(start + i)
        page = requests.get(URL, headers={"User-Agent": "html"})
        soup = BeautifulSoup(page.content, "html.parser")
        logo = soup.find("img")
        printLogo = logo.attrs["src"]
        logoURL = "http://projectabcd.com/" + printLogo
        r = requests.get(logoURL, headers={"User-Agent": "html"}, stream=True)
        if r.status_code == 200:
             with open(basename(printLogo), "wb") as f:
                  r.raw.decode_content = True
                  shutil.copyfileobj(r.raw, f)
        pageInfo = soup.find("div", class_="containerTitle")
        pageInfoImg = soup.find("div", class_="container")
        name = pageInfo.find("h2", class_="headTwo")
        printName = name.text
        image = pageInfoImg.find("div", class_="containerImage")
        img = image.find("image",class_= "image")
        printImage = img.get("src")
        pictureURL = "http://projectabcd.com/" + printImage
        r = requests.get(pictureURL, headers={"User-Agent": "html"}, stream=True)
        if r.status_code == 200:
             with open(basename(printImage), "wb") as f:
                  r.raw.decode_content = True
                  shutil.copyfileobj(r.raw, f)
        element = pageInfo.p
        pagetext = pageInfoImg.find("div", class_="containerText")
        description = pagetext.find("p", class_="words")
        printDescription = description.text
        fact = description.find_next_sibling("p")
        printFact = fact.text
        printWords = printDescription + printFact

        #creates the slide presentation if slide option 1 is choosen
        if(slideOption == 1 ):
             #creates the slides and sets layout preferences
             slide_layout = prs.slide_layouts[6]
             slide = prs.slides.add_slide(slide_layout)
             prs.slide_width = pptx.util.Inches(8)
             prs.slide_height = pptx.util.Inches(11)
             #places the logo on the slide
             logoHolder = slide.shapes.add_picture(basename(printLogo), pptx.util.Inches(7), pptx.util.Inches(0),width=pptx.util.Inches(1), height=pptx.util.Inches(1))
             #places the title on the slide
             titleBox = slide.shapes.add_textbox(pptx.util.Inches(2.5), pptx.util.Inches(.5), width = pptx.util.Inches(3), height = pptx.util.Inches(1))
             titleBoxtf = titleBox.text_frame
             title = titleBoxtf.add_paragraph()
             title.text = printName
             title.font.name = titleFont
             title.font.size = Pt(titleSize)
             #places the picture on the slide
             pictureHolder = prs.slides[i].shapes
             pictureHolder.add_picture(basename(printImage), pptx.util.Inches(2.5),pptx.util.Inches(2),width = pptx.util.Inches(3), height = pptx.util.Inches(4))
             #creates a textbox for the description and fun fact
             contentBox = slide.shapes.add_textbox(pptx.util.Inches(1), pptx.util.Inches(6), width = pptx.util.Inches(6), height = pptx.util.Inches(5))
             contentBoxtf = contentBox.text_frame
             contentBoxtf.word_wrap = True
             descriptionTitle = contentBoxtf.add_paragraph()
             descriptionTitle.font.name = textFont
             descriptionTitle.font.bold = True
             descriptionTitle.font.size = Pt(textSize)
             descriptionTitle.text = "Description: "
             descriptionParagraph = contentBoxtf.add_paragraph()
             descriptionParagraph.font.name = textFont
             descriptionParagraph.font.size = Pt(textSize)
             descriptionParagraph.text = printDescription
             FunFactTitle = contentBoxtf.add_paragraph()
             FunFactTitle.font.name = textFont
             FunFactTitle.font.bold = True
             FunFactTitle.font.size = Pt(textSize)
             FunFactTitle.text = "\nFun Fact:"
             FunFactParagraph = contentBoxtf.add_paragraph()
             FunFactParagraph.font.name = textFont
             FunFactParagraph.font.size = Pt(textSize)
             FunFactParagraph.text = printFact
        #creates the slide presentation if slide option 2 is chosen
        elif(slideOption == 2):
             #creates the slide layout preferences
             slide_layout = prs.slide_layouts[6]
             prs.slide_width = pptx.util.Inches(8)
             prs.slide_height = pptx.util.Inches(11)
             #creates a title page
             if(i == 0):
                  slide1 = prs.slides.add_slide(slide_layout)
                  titleBox = slide1.shapes.add_textbox(pptx.util.Inches(1.5), pptx.util.Inches(4.5),width=pptx.util.Inches(5), height=pptx.util.Inches(2))
                  titleBoxtf = titleBox.text_frame
                  title = titleBoxtf.add_paragraph()
                  title.text = "Project abcd abdul"
                  title.font.size = Pt(60)
                  title.font.name = titleFont
             slide = prs.slides.add_slide(slide_layout)
             #places the picture to cover the whole slide
             pictureHolder = prs.slides[pictureSlide].shapes
             pictureHolder.add_picture(basename(printImage), pptx.util.Inches(0), pptx.util.Inches(0),width=pptx.util.Inches(8), height=pptx.util.Inches(11))
             #creates next slide
             slide2 = prs.slides.add_slide(slide_layout)
             #places the logo on the slide
             logoHolder = slide2.shapes.add_picture(basename(printLogo), pptx.util.Inches(7), pptx.util.Inches(0),width=pptx.util.Inches(1), height=pptx.util.Inches(1))
             #places title on the slide
             titleBox = slide2.shapes.add_textbox(pptx.util.Inches(2), pptx.util.Inches(.5), width=pptx.util.Inches(4),height=pptx.util.Inches(1))
             titleBoxtf = titleBox.text_frame
             title = titleBoxtf.add_paragraph()
             title.text = printName
             title.font.size = Pt(titleSize)
             title.font.name = titleFont
             #creates textbox for description and fun fact
             contentBox = slide2.shapes.add_textbox(pptx.util.Inches(1), pptx.util.Inches(2), width=pptx.util.Inches(6),height=pptx.util.Inches(7))
             contentBoxtf = contentBox.text_frame
             contentBoxtf.word_wrap = True
             descriptionTitle = contentBoxtf.add_paragraph()
             descriptionTitle.font.name = textFont
             descriptionTitle.font.bold = True
             descriptionTitle.font.size = Pt(textSize)
             descriptionTitle.text = "Description: "
             descriptionParagraph = contentBoxtf.add_paragraph()
             descriptionParagraph.font.name = textFont
             descriptionParagraph.font.size = Pt(textSize)
             descriptionParagraph.text = printDescription
             FunFactTitle = contentBoxtf.add_paragraph()
             FunFactTitle.font.bold = True
             FunFactTitle.font.name = textFont
             FunFactTitle.font.size = Pt(textSize)
             FunFactTitle.text = "\nFun Fact:"
             FunFactParagraph = contentBoxtf.add_paragraph()
             FunFactParagraph.font.name = textFont
             FunFactParagraph.font.size = Pt(textSize)
             FunFactParagraph.text = printFact
             pictureSlide = pictureSlide + 2
        #creates the slide presentation if slide option 3 is chosen
        elif (slideOption == 3):
             #creates slide preferences
             slide_layout = prs.slide_layouts[6]
             prs.slide_width = pptx.util.Inches(8)
             prs.slide_height = pptx.util.Inches(11)
             slide2 = prs.slides.add_slide(slide_layout)
             #places picture to cover whole slide
             pictureHolder = prs.slides[pictureSlide].shapes
             pictureHolder.add_picture(basename(printImage), pptx.util.Inches(0), pptx.util.Inches(0),width=pptx.util.Inches(8), height=pptx.util.Inches(11))
             #creates next slide
             slide3 = prs.slides.add_slide(slide_layout)
             #place logo on the slide
             logoHolder = slide3.shapes.add_picture(basename(printLogo), pptx.util.Inches(7), pptx.util.Inches(0),width=pptx.util.Inches(1), height=pptx.util.Inches(1))
             #places the title
             titleBox = slide3.shapes.add_textbox(pptx.util.Inches(4), pptx.util.Inches(.5),width=pptx.util.Inches(4), height=pptx.util.Inches(1))
             titleBoxtf = titleBox.text_frame
             title = titleBoxtf.add_paragraph()
             title.text = printName
             title.font.size = Pt(titleSize)
             title.font.name = titleFont
             #creates textbox for description and fun fact
             contentBox = slide3.shapes.add_textbox(pptx.util.Inches(1), pptx.util.Inches(2), width=pptx.util.Inches(6),height=pptx.util.Inches(7))
             contentBoxtf = contentBox.text_frame
             contentBoxtf.word_wrap = True
             descriptionTitle = contentBoxtf.add_paragraph()
             descriptionTitle.font.name = textFont
             descriptionTitle.font.bold = True
             descriptionTitle.font.size = Pt(textSize)
             descriptionTitle.text = "Description: "
             descriptionParagraph = contentBoxtf.add_paragraph()
             descriptionParagraph.font.name = textFont
             descriptionParagraph.font.size = Pt(textSize)
             descriptionParagraph.text = printDescription
             FunFactTitle = contentBoxtf.add_paragraph()
             FunFactTitle.font.bold = True
             FunFactTitle.font.name = textFont
             FunFactTitle.font.size = Pt(textSize)
             FunFactTitle.text = "\nFun Fact:"
             FunFactParagraph = contentBoxtf.add_paragraph()
             FunFactParagraph.font.name = textFont
             FunFactParagraph.font.size = Pt(textSize)
             FunFactParagraph.text = printFact
             pictureSlide = pictureSlide + 2
   test = "test.pptx"
   prs.save(test)
   return test

#creates and opens the powerpoint presentation
test = buildPresentation()
os.startfile(test)