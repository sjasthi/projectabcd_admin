from pptx import Presentation
import mysql.connector
import os
import collections
import pptx.util
from pptx.util import Pt
from pptx.api import Presentation


conn = mysql.connector.connect(user='root', host='localhost', database='abcd_dress-500')
cursor = conn.cursor()


sql = "SELECT * FROM dresses"
cursor.execute(sql)
result = cursor.fetchall()

def buildPresentation(data):
    start = int(input("Enter the slide to start with: "))
    end = int(input("Enter the slide to end with: "))
    print("Creating PowerPoint slides.")

    with open("preferences.txt", "r") as f:
        slideOption = int(f.readline().split("= ")[1])
        textFont = f.readline().split("= ")[1]
        titleFont = f.readline().split("= ")[1]
        textSize = int(f.readline().split("= ")[1])
        titleSize = int(f.readline().split("= ")[1])

    prs = Presentation()
    presentationLength = end - start + 1
    start = start -1

    for i in range(0,presentationLength):
        
        slide = prs.slides.add_slide(prs.slide_layouts[6]) 
        slide2 = prs.slides.add_slide(prs.slide_layouts[6])
        
        prs.slide_width = pptx.util.Inches(8)
        prs.slide_height = pptx.util.Inches(11)        
        
        contentBox = slide2.shapes.add_textbox(pptx.util.Inches(1), pptx.util.Inches(2), width=pptx.util.Inches(6),height=pptx.util.Inches(7))
        titleBox = slide2.shapes.add_textbox(pptx.util.Inches(4), pptx.util.Inches(.5),width=pptx.util.Inches(4), height=pptx.util.Inches(1))
        
        titleBoxtf = titleBox.text_frame
        title = titleBoxtf.add_paragraph()
        title.font.name = titleFont
        title.font.size = Pt(titleSize)
        title.font.name = titleFont
        title.text = str(data[start + i][1]) 
        
        contentBoxtf = contentBox.text_frame
        contentBoxtf.word_wrap = True
        descriptionTitle = contentBoxtf.add_paragraph()
        descriptionTitle.font.name = textFont
        descriptionTitle.font.bold = True
        descriptionTitle.font.size = Pt(textSize)
        descriptionTitle.text = "Description: "
        descriptionParagraph = contentBoxtf.add_paragraph()
        descriptionParagraph.font.name = textFont
        descriptionParagraph.font.size = Pt(textSize)
        descriptionParagraph.text = str(data[start + i][2])
        FunFactTitle = contentBoxtf.add_paragraph()
        FunFactTitle.font.bold = True
        FunFactTitle.font.name = textFont
        FunFactTitle.font.size = Pt(textSize)
        FunFactTitle.text = "\nFun Fact:"
        FunFactParagraph = contentBoxtf.add_paragraph()
        FunFactParagraph.font.name = textFont
        FunFactParagraph.font.size = Pt(textSize)
        FunFactParagraph.text = str(data[start + i][3])
        
        
        
        image_url = data[start + i][8]  
        if image_url:
            image_path = os.path.basename(image_url)
            slide.shapes.add_picture(image_path, pptx.util.Inches(0), pptx.util.Inches(0),width=pptx.util.Inches(8), height=pptx.util.Inches(11))
        
    test = "dbTest.pptx"
    prs.save(test)
    return test

test = buildPresentation(result)
os.startfile(test)
