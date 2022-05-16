import json
import os
import io
import shutil
from os.path import basename

from bs4 import BeautifulSoup
import requests
import pptx.util
from pptx import Presentation
from pptx.util import Pt
from sqlalchemy import null
import matplotlib.pyplot as plt

#def to allow attribute access to dictionaries
def _init_(self):
    self = dict()

#def to be able to add key value pairs to dictionaries
def add(self, key, value):
    self[key] = value

#initializes dictionaries
categoryDict = {}
typeDict = {}
keywordDict = {}
stateDict = {}
statusDict = {}
descriptionDict = {}

def summarySlide(prs, dictionary, logo):
    # creates a new slide for the summary
    slide_layout = prs.slide_layouts[6]
    slide2 = prs.slides.add_slide(slide_layout)
    prs.slide_width = pptx.util.Inches(8)
    prs.slide_height = pptx.util.Inches(11)
    logoHolder = slide2.shapes.add_picture(basename(logo), pptx.util.Inches(7), pptx.util.Inches(0),width=pptx.util.Inches(1), height=pptx.util.Inches(1))
    contentBox = slide2.shapes.add_textbox(pptx.util.Inches(1), pptx.util.Inches(2), width=pptx.util.Inches(6),height=pptx.util.Inches(5))
    contentBoxtf = contentBox.text_frame
    contentBoxtf.word_wrap = True
    typeParagraph = contentBoxtf.add_paragraph()
    typeParagraph.font.size = Pt(16)
    # creates new table
    x, y, cx, cy = pptx.util.Inches(2), pptx.util.Inches(2), pptx.util.Inches(4), pptx.util.Inches(1.5)
    shape = slide2.shapes.add_table(len(dictionary.keys()), 2, x, y, cx, cy)
    table = shape.table
    count = 0
    # fills table with the summary data
    for i in dictionary:
        if (count == 0):
            cell0 = table.cell(0, 0)
            cell0.text = "Type"
            cell1 = table.cell(0, 1)
            cell1.text = "Count"
            count += 1
        else:
            cell = table.cell(count, 0)
            cell.text = json.dumps(i).strip('"')
            cell2 = table.cell(count, 1)
            cell2.text = json.dumps(dictionary[i])
            count += 1

def summaryGraph(prs, dictionary, logo):
    # creates the slide for the summary graph
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    prs.slide_width = pptx.util.Inches(8)
    prs.slide_height = pptx.util.Inches(11)
    logoHolder = slide.shapes.add_picture(basename(logo), pptx.util.Inches(7), pptx.util.Inches(0),width=pptx.util.Inches(1), height=pptx.util.Inches(1))
    keys = dictionary.keys()
    values = dictionary.values()
    # creates a bar graph using the keys and values as the x and y variables
    plt.bar(keys, values)
    image_stream = io.BytesIO()
    plt.savefig(image_stream)
    x, y, cx, cy = pptx.util.Inches(2), pptx.util.Inches(2), pptx.util.Inches(6), pptx.util.Inches(4.5)
    pic = slide.shapes.add_picture(image_stream, x, y, cx, cy)

#uses the data collected from web scraping to fill the dictionaries for each category
def createDict(item, dictionary):
    if (item.lower() in dictionary):
        dictionary[item.lower()] += 1
    else:
        dictionary[item.lower()] = 1
def buildPresentation():
    print("Creating powerpoint slides.")
    #webscrapes the URL using beautiful soup to get up to date info and pictures
    URL = "https://projectabcd.com/list_dresses.php"
    page = requests.get(URL, headers={"User-Agent": "html"})
    soup = BeautifulSoup(page.content, "html.parser")
    logo = soup.find("img")
    #gets logo image and downloads it onto local machine
    printLogo = logo.attrs["src"]
    logoURL = "http://projectabcd.com/" + printLogo
    r = requests.get(logoURL, headers={"User-Agent": "html"}, stream=True)
    if r.status_code == 200:
         with open(basename(printLogo), "wb") as f:
              r.raw.decode_content = True
              shutil.copyfileobj(r.raw, f)
    section = soup.find("tr")
    #loops through to collect all of the needed information for each item in the list
    #so it can then be categorized in the dictionaries
    while(section.find_next_sibling("tr")!= null):
        next = section.find_next("tr")
        if(next == None):
            break
        id = next.find("td")
        name = id.find_next_sibling("td")
        description = name.find_next_sibling("td")
        #sections the description count into 10 groups by 100s
        if(len(description.text) > 0 and len(description.text) <= 100):
            descriptionType = "0 - 100"
        elif(len(description.text) > 100 and len(description.text) <= 200):
            descriptionType = "101 - 200"
        elif(len(description.text) > 200 and len(description.text) <= 300):
            descriptionType = "201 - 300"
        elif(len(description.text) > 300 and len(description.text) <= 400):
            descriptionType = "301 - 400"
        elif(len(description.text) > 400 and len(description.text) <= 500):
            descriptionType = "401 - 500"
        elif(len(description.text) > 500 and len(description.text) <= 600):
            descriptionType = "501 - 600"
        elif(len(description.text) > 600 and len(description.text) <= 700):
            descriptionType = "601 - 700"
        elif(len(description.text) > 700 and len(description.text) <= 800):
            descriptionType = "701 - 800"
        elif(len(description.text) > 800 and len(description.text) <= 900):
            descriptionType = "801 - 900"
        elif(len(description.text) > 900 and len(description.text) <= 1000):
            descriptionType = "900 - 100"
        createDict(descriptionType, descriptionDict)
        fact = description.find_next_sibling("td")

        category = fact.find_next_sibling("td")
        #creates the dictionary of categories
        if(category is not None and category != ""):
            categories = category.text.split(",")
        for cat in categories:
            createDict(cat, categoryDict)

        type = category.find_next_sibling("td")
        #creates the dictionary of types
        createDict(type.text, typeDict)

        stateName = type.find_next_sibling("td")
        #creates the dictionary of state names
        createDict(stateName.text, stateDict)

        keywords = stateName.find_next_sibling("td")
        #creates the dictionary of keywords
        if (keywords is not None and keywords != ""):
            keywordsList = keywords.text.split(",")
        for key in keywordsList:
            createDict(key, keywordDict)

        status = keywords.find_next_sibling("td")
        #creates the dictionary of statuses
        createDict(status.text, statusDict)
        section = next

    #creates presentation
    prs = Presentation()
    slide_layout = prs.slide_layouts[6]
    #creates all of the summary tables and graphs for each category
    summarySlide(prs, categoryDict, printLogo)
    summaryGraph(prs, categoryDict, printLogo)
    summarySlide(prs, typeDict, printLogo)
    summaryGraph(prs, typeDict, printLogo)
    summarySlide(prs, keywordDict, printLogo)
    summaryGraph(prs, keywordDict, printLogo)
    summarySlide(prs, stateDict, printLogo)
    summaryGraph(prs, statusDict, printLogo)
    summarySlide(prs, statusDict, printLogo)
    summaryGraph(prs, statusDict, printLogo)
    summarySlide(prs, descriptionDict, printLogo)
    summaryGraph(prs, descriptionDict, printLogo)
    summary = "summary.pptx"
    prs.save(summary)
    return summary
summary = buildPresentation()
os.startfile(summary)